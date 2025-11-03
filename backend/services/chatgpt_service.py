from openai import OpenAI
try:
    # Production when imported as backend.services.chatgpt_service
    from ..config.settings import Config
except Exception:
    # Local run from backend/ directory
    from config.settings import Config
from datetime import datetime
from typing import List, Dict, Any
from flask import g

def debug_log(message: str, category: str = "general"):
    """Conditional debug logging based on configuration"""
    if category == "odoo_data" and Config.DEBUG_ODOO_DATA:
        print(f"DEBUG: {message}")
    elif category == "bot_logic" and Config.DEBUG_BOT_LOGIC:
        print(f"DEBUG: {message}")
    elif category == "knowledge_base" and Config.DEBUG_KNOWLEDGE_BASE:
        print(f"DEBUG: {message}")
    elif category == "general" and Config.VERBOSE_LOGS:
        print(f"DEBUG: {message}")
import time
import json
import os
import threading

class ChatGPTService:
    """Service class for handling ChatGPT interactions using the latest available GPT model"""
    
    def __init__(self):
        self.client = None  # Initialize lazily to avoid hanging
        self.model = Config.GPT_MODEL  # Using configured GPT model
        self.fallback_model = getattr(Config, 'GPT_FALLBACK_MODEL', None)
        self.conversation_history = {}  # Store conversation history by thread_id
        self.storage_dir = "conversation_storage"
        self.lock = threading.Lock()  # Thread safety for file operations
        
        # Service dependencies (set via set_services)
        self.timeoff_service = None
        self.session_manager = None
        self.halfday_service = None
        self.reimbursement_service = None
        self.metrics_service = None
        self.leave_balance_service = None

        # Leave types that require the user to pick between full days vs custom hours
        self.leave_mode_config = {
            'Sick Leave': {
                'session_key': 'sick_leave_mode',
                'button_prefix': 'SICK',
                'button_type': 'sick_leave_mode',
                'prompt_label': 'Sick Leave'
            },
            'Unpaid Leave': {
                'session_key': 'unpaid_leave_mode',
                'button_prefix': 'UNPAID',
                'button_type': 'unpaid_leave_mode',
                'prompt_label': 'Unpaid Leave'
            }
        }
        
        # Disable conversation storage
        # if not os.path.exists(self.storage_dir):
        #     os.makedirs(self.storage_dir)

    def get_current_odoo_session(self):
        """Get Odoo session data from Flask's request context (g object)"""
        return getattr(g, 'odoo_session_data', None)

    def set_services(
        self,
        timeoff_service,
        session_manager,
        halfday_service=None,
        reimbursement_service=None,
        metrics_service=None
    ):
        """Wire external services for advanced functionality"""
        self.timeoff_service = timeoff_service
        self.session_manager = session_manager
        self.halfday_service = halfday_service
        self.reimbursement_service = reimbursement_service
        self.metrics_service = metrics_service

    def _resolve_identity(self, employee_data: dict = None) -> Dict[str, str]:
        """Extract tenant/user identifiers for metric logging."""
        tenant_id = None
        tenant_name = None
        user_id = None
        user_name = None
        if isinstance(employee_data, dict):
            try:
                eid = employee_data.get('id')
                if eid is not None:
                    user_id = str(eid)
            except Exception:
                pass
            try:
                company_details = employee_data.get('company_id_details')
                if isinstance(company_details, dict):
                    if company_details.get('id') is not None:
                        tenant_id = str(company_details.get('id'))
                    if company_details.get('name'):
                        tenant_name = company_details.get('name')
                else:
                    raw_company = employee_data.get('company_id')
                    if isinstance(raw_company, (list, tuple)) and raw_company:
                        tenant_id = str(raw_company[0])
                        if len(raw_company) > 1 and raw_company[1]:
                            tenant_name = raw_company[1]
                    elif raw_company is not None:
                        tenant_id = str(raw_company)
                try:
                    name = employee_data.get('name')
                    if name:
                        user_name = str(name)
                except Exception:
                    pass
            except Exception:
                pass
        payload_meta = {}
        if tenant_name:
            payload_meta['tenant_name'] = tenant_name
        return {
            'tenant_id': tenant_id,
            'user_id': user_id,
             'user_name': user_name,
            'meta': payload_meta
        }

    def _record_metric(self, metric_type: str, thread_id: str, payload: Dict[str, Any], employee_data: dict = None) -> bool:
        """Send metric to Supabase if the integration is configured."""
        try:
            if not self.metrics_service:
                return False
            identity = self._resolve_identity(employee_data)
            merged_payload = dict(payload or {})
            meta = identity.get('meta')
            if meta:
                merged_payload.setdefault('context', {}).update(meta)
            success = self.metrics_service.log_metric(
                metric_type,
                thread_id,
                user_id=identity.get('user_id'),
                user_name=identity.get('user_name'),
                tenant_id=identity.get('tenant_id'),
                payload=merged_payload
            )
            return success
        except Exception:
            # Never let metrics interfere with the main flow
            return False

    def _is_timeoff_start_message(self, text):
        """Return True when the user explicitly asks to begin a fresh time-off flow."""
        try:
            phrase = (text or '').strip().lower()
            if not phrase:
                return False
            blockers = {'cancel', 'stop', 'exit', 'quit', 'abort', 'undo', 'no', 'n', 'yes', 'y', 'confirm', 'submit', 'ok', 'sure'}
            if any(tok in phrase for tok in blockers):
                return False
            keywords = ['time off', 'day off', 'leave', 'vacation', 'holiday', 'rest day']
            verbs = ['i want', 'i need', 'i would like', 'request', 'apply', 'book', 'submit', 'take', 'get', 'start', 'begin']
            # Require an explicit verb to treat as a fresh start; plain leave type names (e.g., 'annual leave')
            # should be handled inside the flow as a selection, not as a restart trigger.
            if any(k in phrase for k in keywords) and any(v in phrase for v in verbs):
                return True
            # Allow very short explicit starts only for generic phrases, not specific leave types
            single_intents = {'time off', 'request time off', 'apply for leave', 'submit leave request'}
            return phrase in single_intents
        except Exception:
            return False

    def _is_timeoff_continuation_message(self, text):
        """Return True when the message looks like a reply within an existing time-off flow."""
        try:
            phrase = (text or '').strip().lower()
            if not phrase:
                return False
            start_keywords = ['time off', 'leave', 'annual leave', 'sick leave', 'request time off']
            if any(k in phrase for k in start_keywords):
                return False
            continuation_tokens = {
                'yes', 'y', 'no', 'n', 'submit', 'confirm', 'cancel', 'stop',
                'exit', 'quit', 'ok', 'sure'
            }
            if phrase in continuation_tokens:
                return True
            if any(tok in phrase for tok in ['hour_from=', 'hour_to=']):
                return True
            if any(term in phrase for term in ['annual', 'sick', 'custom hours', 'work from home']):
                return True
            if phrase.isdigit() and len(phrase) <= 2:
                return True
            import re as _re
            if _re.search(r"\d{1,2}[/-]\d{1,2}(?:[/-]\d{2,4})?", phrase):
                return True
            if any(marker in phrase for marker in [' to ', ' until ', ' till ', '-', ' next ', ' tomorrow']):
                return True
            return False
        except Exception:
            return False

    def _reset_timeoff_sessions(self, thread_id, reason='', employee_data=None):
        """
        Force-remove lingering time-off sessions for the CURRENT USER only.

        CRITICAL: This now only clears sessions belonging to the same employee
        to support concurrent timeoff requests from multiple users.

        Args:
            thread_id: Current thread ID to clear
            reason: Reason for clearing
            employee_data: Current user's employee data (to identify their sessions)
        """
        try:
            # Get current user's employee ID for filtering
            current_employee_id = None
            if employee_data and isinstance(employee_data, dict):
                current_employee_id = employee_data.get('id')

            # Clear the current thread's session
            if thread_id:
                try:
                    if reason:
                        self.session_manager.cancel_session(thread_id, reason)
                except Exception:
                    pass
                try:
                    self.session_manager.clear_session(thread_id)
                except Exception:
                    pass

            # CRITICAL CHANGE: Only clear OTHER sessions belonging to the SAME employee
            # This prevents clearing sessions from other concurrent users
            if current_employee_id:
                debug_log(f"Clearing other timeoff sessions for employee {current_employee_id}", "bot_logic")
                try:
                    for tid, sess in list(getattr(self.session_manager, 'sessions', {}).items()):
                        if tid == thread_id:
                            continue
                        if isinstance(sess, dict) and sess.get('type') == 'timeoff':
                            # Check if this session belongs to the same employee
                            sess_employee_id = None
                            try:
                                sess_employee_data = sess.get('data', {}).get('employee_data', {})
                                if not sess_employee_data:
                                    sess_employee_data = sess.get('employee_data', {})
                                sess_employee_id = sess_employee_data.get('id')
                            except Exception:
                                pass

                            # Only clear if it's the same employee
                            if sess_employee_id == current_employee_id:
                                debug_log(f"Clearing session {tid} for same employee {current_employee_id}", "bot_logic")
                                try:
                                    self.session_manager.clear_session(tid)
                                except Exception:
                                    continue
                            else:
                                debug_log(f"Keeping session {tid} - belongs to different employee {sess_employee_id}", "bot_logic")
                except Exception as e:
                    debug_log(f"Error in session filtering: {e}", "general")
                    pass

                try:
                    for tid, sess in self.session_manager.find_active_timeoff_sessions():
                        if thread_id and tid == thread_id:
                            continue
                        # Check if this session belongs to the same employee
                        sess_employee_id = None
                        try:
                            sess_employee_data = sess.get('data', {}).get('employee_data', {})
                            if not sess_employee_data:
                                sess_employee_data = sess.get('employee_data', {})
                            sess_employee_id = sess_employee_data.get('id')
                        except Exception:
                            pass

                        # Only clear if it's the same employee
                        if sess_employee_id == current_employee_id:
                            debug_log(f"Clearing active session {tid} for same employee {current_employee_id}", "bot_logic")
                            try:
                                self.session_manager.clear_session(tid)
                            except Exception:
                                continue
                        else:
                            debug_log(f"Keeping active session {tid} - belongs to different employee {sess_employee_id}", "bot_logic")
                except Exception as e:
                    debug_log(f"Error in active session filtering: {e}", "general")
                    pass
            else:
                debug_log("No employee_id provided - skipping multi-user session cleanup for safety", "bot_logic")
        except Exception as e:
            debug_log(f"Error in _reset_timeoff_sessions: {e}", "general")
            pass

    def _restart_timeoff_flow(self, message: str, thread_id: str, employee_data: dict, reason: str = '') -> dict:
        """Helper to force-reset and start a fresh time-off flow in the middle of another step."""
        extracted_payload = {}
        try:
            detected = self.timeoff_service.detect_timeoff_intent(message)
            if detected[0]:
                extracted_payload = detected[2] or {}
        except Exception:
            extracted_payload = {}

        self._reset_timeoff_sessions(thread_id, reason or 'User requested new time-off flow', employee_data)
        return self._start_timeoff_session(message, thread_id, extracted_payload, employee_data)

    def _persist_timeoff_context(self, thread_id: str, session: dict, **fields) -> None:
        """Persist important time-off context (leave type, dates, hours) for later steps).

        Guard against overwriting previously captured values with empty strings or None.
        """
        if not thread_id or not fields:
            return

        try:
            existing_data = {}
            if isinstance(session, dict):
                existing_data = session.get('data', {}) or {}

            context = dict(existing_data.get('timeoff_context', {}))
            for key, value in fields.items():
                # Skip empty-string values to avoid erasing prior state
                if isinstance(value, str) and value.strip() == "":
                    continue
                if value is not None:
                    context[key] = value

            update_payload = {}
            for key, value in fields.items():
                # Skip empty-string values to avoid erasing prior state
                if isinstance(value, str) and value.strip() == "":
                    continue
                if value is not None:
                    update_payload[key] = value

            update_payload['data'] = {**existing_data, 'timeoff_context': context}
            self.session_manager.update_session(thread_id, update_payload)
        except Exception as persist_error:
            debug_log(f"Failed to persist time-off context: {persist_error}", "general")
    
    # -------------------- Leave mode helpers --------------------
    def _mode_config_for_leave(self, leave_name: str) -> dict:
        if not leave_name:
            return {}
        return self.leave_mode_config.get(leave_name.strip(), {})

    def _leave_mode_key(self, leave_name: str) -> str:
        cfg = self._mode_config_for_leave(leave_name)
        return cfg.get('session_key')

    def _get_leave_mode(self, session: dict, leave_name: str) -> str:
        key = self._leave_mode_key(leave_name)
        if not key:
            return None
        try:
            session_data = session.get('data', {}) if isinstance(session, dict) else {}
            ctx = {}
            if isinstance(session_data, dict):
                ctx = session_data.get('timeoff_context', {}) or {}
            return (
                ctx.get(key)
                or session_data.get(key)
                or session.get(key)
            )
        except Exception:
            return None

    def _store_leave_mode(self, thread_id: str, session: dict, leave_name: str, mode: str) -> None:
        key = self._leave_mode_key(leave_name)
        if not key or mode is None:
            return
        try:
            self.session_manager.update_session(thread_id, {key: mode})
        except Exception:
            pass
        try:
            self._persist_timeoff_context(thread_id, session, **{key: mode})
        except Exception:
            pass

    def _prompt_leave_mode(self, thread_id: str, leave_name: str, prefix_note: str = "") -> dict:
        cfg = self._mode_config_for_leave(leave_name)
        label = cfg.get('prompt_label') or leave_name or 'this leave type'
        prefix = cfg.get('button_prefix') or label.upper().replace(' ', '_')
        button_type = cfg.get('button_type') or f"{prefix.lower()}_mode"
        text = ((prefix_note + "\n\n") if prefix_note else "")
        text += f"For {label}, do you want to take Full Days or Custom Hours?"
        buttons = [
            {'text': 'Full Days', 'value': f"{prefix}_FULL_DAYS", 'type': button_type},
            {'text': 'Custom Hours', 'value': f"{prefix}_CUSTOM_HOURS", 'type': button_type}
        ]
        return self._create_response_with_choice_buttons(text, thread_id, buttons)

    def _is_custom_hours_mode(self, session: dict, leave_name: str) -> bool:
        return self._get_leave_mode(session, leave_name) == 'custom_hours'

    def _get_client(self):
        """Get OpenAI client, initializing lazily if needed"""
        if self.client is None:
            try:
                self.client = OpenAI(api_key=Config.OPENAI_API_KEY)
            except Exception as e:
                debug_log(f"Failed to initialize OpenAI client: {e}", "general")
                return None
        return self.client
    
    def _get_storage_file(self, thread_id):
        """Get the storage file path for a thread"""
        return os.path.join(self.storage_dir, f"thread_{thread_id}.json")

    def _get_summary_file(self, thread_id):
        return os.path.join(self.storage_dir, f"thread_{thread_id}_summary.json")

    def _load_summary(self, thread_id):
        try:
            path = self._get_summary_file(thread_id)
            if os.path.exists(path):
                with open(path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return data.get('summary', '')
        except Exception as e:
            debug_log(f"Error loading summary: {e}", "general")
        return ''

    def _save_summary(self, thread_id, summary_text):
        try:
            path = self._get_summary_file(thread_id)
            with open(path, 'w', encoding='utf-8') as f:
                json.dump({'summary': summary_text or ''}, f, ensure_ascii=False, indent=2)
        except Exception as e:
            debug_log(f"Error saving summary: {e}", "general")
    
    def _load_conversation_history(self, thread_id):
        """Load conversation history from persistent storage - DISABLED"""
        # Storage disabled - return empty history
        return []
    
    def _save_conversation_history(self, thread_id, history):
        """Save conversation history to persistent storage - DISABLED"""
        # Storage disabled - no longer saving conversation history
        pass
    
    def get_response(self, message, thread_id, employee_data=None):
        """Get response from configured GPT model using Chat Completions API"""
        try:
            # Check for active time-off session or detect new time-off intent
            if self.timeoff_service and self.session_manager:
                try:
                    if self._is_timeoff_start_message(message):
                        detected = (False, 0.0, {})
                        try:
                            detected = self.timeoff_service.detect_timeoff_intent(message)
                        except Exception:
                            detected = (False, 0.0, {})
                        extracted_payload = detected[2] if detected[0] else {}
                        self._reset_timeoff_sessions(thread_id, 'User requested new time-off flow', employee_data)
                        return self._start_timeoff_session(message, thread_id, extracted_payload, employee_data)

                    try:
                        active_for_thread = self.session_manager.get_active_session(thread_id) if thread_id else None
                    except Exception:
                        active_for_thread = None
                    # Do NOT rebind across threads to avoid state leakage between flows
                    if not active_for_thread and self._is_timeoff_continuation_message(message):
                        try:
                            active_list = self.session_manager.find_active_timeoff_sessions()
                            if isinstance(active_list, list) and len(active_list) == 1:
                                rebound_thread_id, rebound_session = active_list[0]
                                thread_id = thread_id or rebound_thread_id
                                debug_log(f"Rebinding to active time-off session thread: {thread_id}", "bot_logic")
                        except Exception:
                            pass

                    timeoff_response = self._handle_timeoff_flow(message, thread_id, employee_data)
                    if timeoff_response:
                        debug_log(f"Time-off flow handled, returning response", "bot_logic")
                        return timeoff_response
                    debug_log(f"Time-off flow did not handle message, continuing to normal flow", "bot_logic")
                except Exception as timeoff_error:
                    debug_log(f"Error in time-off flow: {timeoff_error}", "general")
                    import traceback
                    traceback.print_exc()
                    
                    # Clear any sessions that might be stuck
                    if thread_id:
                        self.session_manager.clear_session(thread_id)
                        
                    # Return a helpful error message instead of hanging
                    return {
                        'message': "I tried to help with your time-off request but encountered a technical issue. Please try again or contact HR for assistance. Error details have been logged for troubleshooting.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }
            
            # Check for active reimbursement session or detect new reimbursement intent
            if self.reimbursement_service and self.session_manager:
                try:
                    debug_log(f"Checking reimbursement flow for message: {message[:50]}...", "bot_logic")
                    reimbursement_response = self.reimbursement_service.handle_flow(message, thread_id, employee_data)
                    if reimbursement_response:
                        debug_log(f"Reimbursement flow handled, returning response", "bot_logic")
                        return reimbursement_response
                    debug_log(f"Reimbursement flow did not handle message, continuing to normal flow", "bot_logic")
                except Exception as reimbursement_error:
                    debug_log(f"Error in reimbursement flow: {reimbursement_error}", "general")
                    import traceback
                    traceback.print_exc()
                    
                    # Clear any sessions that might be stuck
                    if thread_id:
                        self.session_manager.clear_session(thread_id)
                        
                    # Return a helpful error message instead of hanging
                    return {
                        'message': "I tried to help with your reimbursement request but encountered a technical issue. Please try again or contact HR for assistance. Error details have been logged for troubleshooting.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }
            
            # Load conversation history and apply cap/condense strategy
            full_history = self._load_conversation_history(thread_id)
            max_history = getattr(Config, 'MAX_HISTORY_MESSAGES', 20)
            context_limit = getattr(Config, 'HISTORY_CONTEXT_LIMIT', 150000)
            include_condensed = getattr(Config, 'INCLUDE_CONDENSED_HISTORY', True)

            # Split into older and recent
            recent_history = full_history[-max_history:] if full_history else []
            older_history = full_history[:-max_history] if full_history and len(full_history) > max_history else []

            condensed_context = ""
            if older_history:
                # Condense older messages into a compact context paragraph
                try:
                    pairs = []
                    for h in older_history:
                        role = h.get('role')
                        content = h.get('content', '')
                        if role in ('user', 'assistant') and content:
                            tag = 'U' if role == 'user' else 'A'
                            pairs.append(f"[{tag}] {content}")
                    condensed_context = "\n".join(pairs)
                    if len(condensed_context) > context_limit:
                        condensed_context = condensed_context[:context_limit] + "\n..."
                except Exception:
                    condensed_context = ""
            
            # Prepare system message with Nasma rules and facts memory
            system_message = """You are Nasma, a helpful, precise assistant for PrezLab. 

When providing information about company policies, procedures, or guidelines, always give comprehensive, detailed explanations. Break down complex topics into clear sections and provide specific details, examples, and step-by-step processes when applicable.

For policy-related questions, include:
- Complete policy details and requirements
- Specific procedures and processes
- Applicable rates, calculations, or formulas
- Country-specific variations (Jordan, UAE, KSA)
- Approval processes and workflows
- Examples and scenarios when helpful

Be thorough and informative while maintaining clarity and accuracy."""

            # Facts memory (from Odoo + chat-derived preferences placeholder)
            facts_lines = []
            if employee_data:
                facts_lines.append(f"Name: {employee_data.get('name', 'Unknown')}")
                facts_lines.append(f"Job Title: {employee_data.get('job_title', 'Unknown')}")
                dept = employee_data.get('department_id_details', {})
                facts_lines.append(f"Department: {dept.get('name') if isinstance(dept, dict) else 'Unknown'}")
                manager = employee_data.get('parent_id_details', {})
                facts_lines.append(f"Manager: {manager.get('name') if isinstance(manager, dict) else 'Unknown'}")
                facts_lines.append(f"Time Zone: {employee_data.get('tz', 'Unknown')}")
                company = employee_data.get('company_id_details', {})
                facts_lines.append(f"Company: {company.get('name') if isinstance(company, dict) else 'Prezlab'}")

            facts_block = "\n".join(["Facts:"] + facts_lines) if facts_lines else ""

            # Rolling conversation summary
            rolling_summary = self._load_summary(thread_id)
            
            # Build messages array per requested layout
            messages = [{"role": "system", "content": system_message}]

            # Simple KB: read and inject as a single system message
            try:
                if False:  # Temporarily disable KB to fix hanging issue
                    kb_dir = os.path.normpath(os.path.join(os.path.dirname(__file__), '..', 'knowledge_base'))
                    max_chars = int(getattr(Config, 'KB_MAX_CHARS', 2000))
                    debug_log(f"Knowledge base enabled - checking directory: {kb_dir}", "knowledge_base")
                    debug_log(f"Max characters limit: {max_chars}", "knowledge_base")
                    
                    if os.path.isdir(kb_dir):
                        kb_texts = []
                        kb_used_files = []
                        all_files = os.listdir(kb_dir)
                        debug_log(f"Found {len(all_files)} files in knowledge base: {all_files}", "knowledge_base")
                        
                        for fname in all_files:
                            lower = fname.lower()
                            path = os.path.join(kb_dir, fname)
                            try:
                                if lower.endswith(('.md', '.txt')):
                                    debug_log(f"Processing text file: {fname}", "knowledge_base")
                                    with open(path, 'r', encoding='utf-8') as f:
                                        content = f.read()
                                        kb_texts.append(f"# {fname}\n" + content)
                                        kb_used_files.append(fname)
                                        debug_log(f"Loaded {fname}: {len(content)} characters", "knowledge_base")
                                elif lower.endswith('.docx'):
                                    try:
                                        debug_log(f"Processing DOCX file: {fname}", "knowledge_base")
                                        from docx import Document
                                        doc = Document(path)
                                        parts = [p.text for p in doc.paragraphs if p.text]
                                        content = "\n".join(parts)
                                        kb_texts.append(f"# {fname}\n" + content)
                                        kb_used_files.append(fname)
                                        debug_log(f"Loaded {fname}: {len(content)} characters", "knowledge_base")
                                    except Exception as de:
                                        debug_log(f"DOCX parse error for {fname}: {de}", "knowledge_base")
                                elif lower.endswith('.pptx'):
                                    try:
                                        debug_log(f"Processing PPTX file: {fname}", "knowledge_base")
                                        from pptx import Presentation
                                        prs = Presentation(path)
                                        slides_text = []
                                        for s in prs.slides:
                                            buf = []
                                            for shp in s.shapes:
                                                if hasattr(shp, 'text') and shp.text:
                                                    buf.append(shp.text)
                                            if buf:
                                                slides_text.append("\n".join(buf))
                                        content = "\n\n".join(slides_text)
                                        kb_texts.append(f"# {fname}\n" + content)
                                        kb_used_files.append(fname)
                                        debug_log(f"Loaded {fname}: {len(content)} characters", "knowledge_base")
                                    except Exception as pe:
                                        debug_log(f"PPTX parse error for {fname}: {pe}", "knowledge_base")
                            except Exception as fe:
                                debug_log(f"KB read error for {fname}: {fe}", "knowledge_base")
                        
                        if kb_texts:
                            # Prioritize important files when truncating
                            priority_files = ['overtime', 'policy', 'leave', 'holiday']
                            priority_texts = []
                            regular_texts = []
                            
                            for text in kb_texts:
                                is_priority = any(priority in text.lower() for priority in priority_files)
                                if is_priority:
                                    priority_texts.append(text)
                                else:
                                    regular_texts.append(text)
                            
                            # Combine with priority files first
                            ordered_texts = priority_texts + regular_texts
                            kb_blob = "\n\n".join(ordered_texts)
                            debug_log(f"Total knowledge base content: {len(kb_blob)} characters", "knowledge_base")
                            debug_log(f"Priority files found: {len(priority_texts)}", "knowledge_base")
                            
                            if len(kb_blob) > max_chars:
                                # Try to keep priority content
                                if priority_texts:
                                    priority_blob = "\n\n".join(priority_texts)
                                    remaining_chars = max_chars - len(priority_blob) - 100  # Reserve space
                                    if remaining_chars > 0:
                                        regular_blob = "\n\n".join(regular_texts)
                                        if len(regular_blob) > remaining_chars:
                                            regular_blob = regular_blob[:remaining_chars] + "\n..."
                                        kb_blob = priority_blob + "\n\n" + regular_blob
                                    else:
                                        kb_blob = priority_blob[:max_chars-3] + "..."
                                else:
                                    kb_blob = kb_blob[:max_chars] + "\n..."
                                
                                debug_log(f"KB content truncated to {len(kb_blob)} characters (prioritized)", "knowledge_base")
                                debug_log(f"KB files loaded (truncated): {kb_used_files}", "knowledge_base")
                            else:
                                debug_log(f"KB files loaded (full): {kb_used_files}", "knowledge_base")
                            
                            debug_log(f"Adding knowledge base to ChatGPT context with {len(kb_blob)} characters", "knowledge_base")
                            messages.append({"role": "system", "content": "KB:\n" + kb_blob})
                        else:
                            debug_log("No knowledge base content loaded", "knowledge_base")
                    else:
                        debug_log(f"KB directory not found: {kb_dir}", "knowledge_base")
                else:
                    debug_log("Knowledge base is disabled in configuration", "knowledge_base")
            except Exception as kb_e:
                debug_log(f"KB injection failed: {kb_e}", "knowledge_base")

            if facts_block:
                messages.append({"role": "system", "content": facts_block})

            if rolling_summary:
                messages.append({"role": "system", "content": "Conversation summary:\n" + rolling_summary})

            # Add only the recent capped history
            for msg in recent_history:
                messages.append({"role": msg["role"], "content": msg["content"]})
            
            # Add current user message
            messages.append({"role": "user", "content": message})
            
            # Build request args depending on model capability
            def _build_chat_args(model_name: str, temperature: float = 1):
                base_args = {
                    'model': model_name,
                    'messages': messages,
                    'temperature': temperature,
                    'top_p': 1.0,
                    'frequency_penalty': 0.0,
                    'presence_penalty': 0.0,
                    'response_format': {'type': 'text'}
                }
                # gpt-5 uses max_completion_tokens instead of max_tokens per error message
                if model_name.startswith('gpt-5'):
                    base_args['max_completion_tokens'] = 2000
                else:
                    base_args['max_tokens'] = 2000
                return base_args

            # Make API call to the configured model with optional fallback
            try:
                client = self._get_client()
                if client is None:
                    return {
                        'message': "I'm experiencing technical difficulties. Please try again later.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 0.0,
                        'model_used': self.model
                    }
                response = client.chat.completions.create(**_build_chat_args(self.model), timeout=30)
            except Exception as primary_error:
                if self.fallback_model:
                    client = self._get_client()
                    if client is None:
                        return {
                            'message': "I'm experiencing technical difficulties. Please try again later.",
                            'thread_id': thread_id,
                            'source': self.model,
                            'confidence_score': 0.0,
                            'model_used': self.model
                        }
                    response = client.chat.completions.create(**_build_chat_args(self.fallback_model))
                    # Update reported model
                    self.model = self.fallback_model
                else:
                    raise

            # Extract response (guard against empty content)
            response_text = getattr(response.choices[0].message, 'content', None)
            if response_text is None or (isinstance(response_text, str) and response_text.strip() == ""):
                # Auto-retry once with safer prompt and lower temperature
                retry_messages = list(messages)
                retry_messages.insert(0, {"role": "system", "content": "Reply in plain text only. If any part seems unsafe or unclear, provide a brief safe explanation instead of returning nothing."})
                try:
                    client = self._get_client()
                    if client is None:
                        return {
                            'message': "I'm experiencing technical difficulties. Please try again later.",
                            'thread_id': thread_id,
                            'source': self.model,
                            'confidence_score': 0.0,
                            'model_used': self.model
                        }
                    retry_response = client.chat.completions.create(**_build_chat_args(self.model, temperature=0.3) | {'messages': retry_messages})
                    response_text = getattr(retry_response.choices[0].message, 'content', None)
                except Exception:
                    pass
                
                # Final fallback text if still empty
                if response_text is None or (isinstance(response_text, str) and response_text.strip() == ""):
                    response_text = "Sorry, I couldn't generate a response just now. Please try rephrasing or asking again."
            
            # Update conversation history (store full turns; condensing happens on send)
            updated_history = list(full_history) if full_history else []
            updated_history.append({"role": "user", "content": message})
            updated_history.append({"role": "assistant", "content": response_text})

            # Update rolling summary (simple heuristic: keep last 1000 chars combining recent messages)
            try:
                summary_source = rolling_summary + "\n" + "\n".join(
                    [f"U: {m['content']}" for m in recent_history if m.get('role') == 'user'] +
                    [f"A: {m['content']}" for m in recent_history if m.get('role') == 'assistant'] +
                    [f"A: {response_text}"]
                )
                new_summary = summary_source[-1000:]
                self._save_summary(thread_id, new_summary)
            except Exception as e:
                debug_log(f"Failed to update rolling summary: {e}", "general")
            
            # Save updated conversation history
            self._save_conversation_history(thread_id, updated_history)
                    
            return {
                'message': response_text,
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 1.0,
                'model_used': self.model
            }
                
        except Exception as e:
            error_msg = f"Error communicating with {self.model}: {str(e)}"
            debug_log(f"{error_msg}", "general")
            return {
                'message': error_msg,
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 0.0,
                'error': True
            }
    
    def clear_conversation_history(self, thread_id):
        """Clear conversation history for a specific thread"""
        try:
            with self.lock:
                storage_file = self._get_storage_file(thread_id)
                if os.path.exists(storage_file):
                    os.remove(storage_file)
                    debug_log(f"Cleared conversation history for thread: {thread_id}", "bot_logic")
                    return True
                return False
        except Exception as e:
            debug_log(f"Error clearing conversation history: {e}", "general")
            return False
    
    def get_conversation_history(self, thread_id):
        """Get conversation history for a specific thread"""
        return self._load_conversation_history(thread_id)
    
    def get_model_info(self):
        """Get information about the current model being used"""
        return {
            'model': self.model,
            'description': f'{self.model} - Latest available OpenAI model with enhanced capabilities',
            'features': ['Text generation', 'Context understanding', 'Employee personalization', 'Conversation history']
        }
    
    def _handle_timeoff_flow(self, message: str, thread_id: str, employee_data: dict) -> dict:
        """Handle time-off request flow with session management"""
        try:
            debug_log(f"Starting time-off flow check...", "bot_logic")

            # Check for active session using thread_id if provided
            active_session = None
            if thread_id:
                active_session = self.session_manager.get_session(thread_id)
                debug_log(f"Active session check for thread_id {thread_id}: {active_session is not None}", "bot_logic")

            # Pre-calc common intent flags so we can reuse them safely
            try:
                start_phrase = self._is_timeoff_start_message(message)
            except Exception:
                start_phrase = False
            try:
                continuation_phrase = self._is_timeoff_continuation_message(message)
            except Exception:
                continuation_phrase = False
            has_timeoff_session = bool(active_session and active_session.get('type') == 'timeoff')

            # Detect new time-off intent first
            is_timeoff, confidence, extracted_data = self.timeoff_service.detect_timeoff_intent(message)
            debug_log(f"Time-off detection complete: is_timeoff={is_timeoff}, confidence={confidence}", "bot_logic")

            wants_timeoff = (
                start_phrase or
                (continuation_phrase and has_timeoff_session) or
                (is_timeoff and confidence >= 0.3) or
                has_timeoff_session
            )

            # Enforce single active flow per thread: block if another flow is active only when user is interacting with time-off
            if wants_timeoff:
                try:
                    active_any = self.session_manager.get_active_session(thread_id) if thread_id else None
                    if active_any and active_any.get('type') not in (None, 'timeoff') and active_any.get('state') in ['started', 'active']:
                        other = active_any.get('type', 'another')
                        return {
                            'message': f"You're currently in an active {other} request. Please complete it or type 'cancel' before starting a new time-off request.",
                            'thread_id': thread_id,
                            'source': self.model,
                            'confidence_score': 1.0,
                            'session_handled': True
                        }
                except Exception:
                    pass

            # If the user explicitly restarts while a session exists, wipe and start fresh
            if start_phrase:
                debug_log("Restart phrase detected during active flow; resetting session.", "bot_logic")
                self._reset_timeoff_sessions(thread_id, 'User restarted time-off flow mid-session', employee_data)
                return self._start_timeoff_session(message, thread_id, extracted_data if is_timeoff else {}, employee_data)

            # If this is a new time-off request (high confidence)
            # BUT: Skip restart if this is a button payload (SICK_FULL_DAYS, SICK_CUSTOM_HOURS, UNPAID_FULL_DAYS, etc)
            is_button_payload = message.strip().upper() in [
                'SICK_FULL_DAYS', 'SICK_CUSTOM_HOURS',
                'UNPAID_FULL_DAYS', 'UNPAID_CUSTOM_HOURS',
                'YES', 'NO', 'CONFIRM', 'CANCEL'
            ]
            if is_timeoff and confidence >= 0.7 and not is_button_payload:
                # Always start a fresh time-off flow on explicit start phrases to avoid bleeding states
                debug_log(f"High confidence time-off intent detected ({confidence:.2f}); forcing a clean start.", "bot_logic")
                try:
                    self._reset_timeoff_sessions(thread_id, 'New time-off request detected - force fresh session', employee_data)
                except Exception:
                    try:
                        if active_session:
                            self.session_manager.clear_session(thread_id)
                    except Exception:
                        pass
                # Validate that we have employee data before starting session
                if not employee_data or not isinstance(employee_data, dict) or not employee_data.get('id'):
                    debug_log(f"Invalid employee data for time-off request: {employee_data}", "bot_logic")
                    return {
                        'message': "I'd like to help with your time-off request, but I need to verify your employee information first. Please try logging out and logging back in, or contact HR for assistance.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }

                debug_log(f"Time-off intent detected with confidence {confidence:.2f}, starting session...", "bot_logic")
                timeoff_result = self._start_timeoff_session(message, thread_id, extracted_data, employee_data)
                debug_log(f"Time-off session start result: {timeoff_result is not None}", "bot_logic")
                if timeoff_result:
                    debug_log(f"Time-off session returned result with message: {timeoff_result.get('message', 'NO MESSAGE')[:50]}...", "bot_logic")
                else:
                    debug_log(f"Time-off session returned None - this is the problem!", "bot_logic")
                return timeoff_result

            # If we have an active session, continue it
            if active_session and active_session.get('type') == 'timeoff':
                # Validate session state - accept both 'started' and 'active' as valid
                session_state = active_session.get('state', 'unknown')
                if session_state in ['completed', 'cancelled']:
                    debug_log(f"Session is {session_state}, clearing and returning None", "bot_logic")
                    self.session_manager.clear_session(thread_id)
                    return None
                elif session_state in ['started', 'active']:
                    debug_log(f"Found active time-off session in step {active_session.get('step', 'unknown')}, continuing...", "bot_logic")
                    return self._continue_timeoff_session(message, thread_id, active_session, employee_data)

            # Lower confidence time-off detection (0.3 - 0.7)
            if is_timeoff and confidence >= 0.3:
                # Validate that we have employee data before starting session
                if not employee_data or not isinstance(employee_data, dict) or not employee_data.get('id'):
                    debug_log(f"Invalid employee data for time-off request: {employee_data}", "bot_logic")
                    return {
                        'message': "I'd like to help with your time-off request, but I need to verify your employee information first. Please try logging out and logging back in, or contact HR for assistance.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }

                debug_log(f"Time-off intent detected with confidence {confidence:.2f}, starting session...", "bot_logic")
                timeoff_result = self._start_timeoff_session(message, thread_id, extracted_data, employee_data)
                debug_log(f"Time-off session start result: {timeoff_result is not None}", "bot_logic")
                if timeoff_result:
                    debug_log(f"Time-off session returned result with message: {timeoff_result.get('message', 'NO MESSAGE')[:50]}...", "bot_logic")
                else:
                    debug_log(f"Time-off session returned None - this is the problem!", "bot_logic")
                return timeoff_result

            debug_log(f"No time-off intent detected (confidence: {confidence:.2f}), returning None", "bot_logic")
            return None

        except Exception as e:
            debug_log(f"Error in time-off flow: {e}", "general")
            import traceback
            traceback.print_exc()
            # Clear any existing session that might be causing issues
            if thread_id:
                self.session_manager.clear_session(thread_id)
            return None
    
    def _start_timeoff_session(self, message: str, thread_id: str, extracted_data: dict, employee_data: dict) -> dict:
        """Start a new time-off request session"""
        try:
            # Validate thread_id
            if not thread_id:
                debug_log("Warning: thread_id is None or empty, generating fallback ID", "bot_logic")
                import time
                thread_id = f"timeoff_{int(time.time())}"
                debug_log(f"Generated fallback thread_id: {thread_id}", "bot_logic")
            
            debug_log(f"Starting time-off session with thread_id: {thread_id}", "bot_logic")
            
            # Clear any existing session + summary first to prevent conflicts
            try:
                self._reset_timeoff_flow_state(thread_id)
            except Exception:
                self.session_manager.clear_session(thread_id)
            
            # Start session
            session_data = {
                'extracted_data': extracted_data,
                'employee_data': employee_data,
                'timeoff_context': {
                    'employee_data': employee_data or {}
                }
            }
            session = self.session_manager.start_session(thread_id, 'timeoff', session_data)

            try:
                self._persist_timeoff_context(thread_id, session, employee_data=employee_data or {})
            except Exception:
                pass
            
            # Get available leave types from Odoo
            success, leave_types = self.timeoff_service.get_leave_types()
            
            if not success:
                # Handle Odoo connection issues gracefully
                debug_log(f"Failed to fetch leave types from Odoo: {leave_types}", "bot_logic")
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I'd like to help with your time-off request, but I'm unable to connect to the HR system right now. Please try again later or contact HR directly for assistance.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Validate and clean leave types data
            if not leave_types or not isinstance(leave_types, list):
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I'm having trouble accessing the available leave types. Please contact HR for assistance with your time-off request.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Clean corrupted leave types data
            debug_log(f"Cleaning {len(leave_types)} leave types for corruption...", "bot_logic")
            clean_leave_types = []
            seen_names = set()
            
            for i, lt in enumerate(leave_types):
                try:
                    # Skip if not a proper dictionary
                    if not isinstance(lt, dict):
                        debug_log(f"Skipping non-dict leave type at index {i}: {type(lt)}", "bot_logic")
                        continue
                    
                    # Must have required fields
                    if 'id' not in lt or 'name' not in lt:
                        debug_log(f"Skipping incomplete leave type at index {i}: {lt}", "bot_logic")
                        continue
                    
                    # Skip duplicates by name
                    name = lt.get('name', '').strip()
                    if not name or name in seen_names:
                        debug_log(f"Skipping duplicate/empty name at index {i}: {name}", "bot_logic")
                        continue
                    
                    # Create clean entry with only essential fields
                    clean_entry = {
                        'id': lt.get('id'),
                        'name': name,
                        'active': lt.get('active', True)
                    }
                    
                    clean_leave_types.append(clean_entry)
                    seen_names.add(name)
                    
                except Exception as clean_error:
                    debug_log(f"Error cleaning leave type {i}: {clean_error}", "bot_logic")
                    continue
            
            debug_log(f"Cleaned leave types: {len(clean_leave_types)} valid entries", "bot_logic")
            
            if not clean_leave_types:
                debug_log(f"No valid leave types after cleaning", "bot_logic")
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "The leave types data from HR system appears to be corrupted. Please contact HR directly for assistance with your time-off request.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Use cleaned data
            leave_types = clean_leave_types

            # Inject Half Days option by replacing Unpaid Leave when applicable
            try:
                if self.halfday_service:
                    leave_types = self.halfday_service.replace_unpaid_with_halfdays(leave_types)
                    debug_log(f"HalfDay: Post-replacement leave types count: {len(leave_types)}", "bot_logic")
            except Exception as hd_e:
                debug_log(f"HalfDay replacement failed: {hd_e}", "general")
            
            # Update session with leave types - with error handling
            try:
                debug_log(f"Updating session with {len(leave_types)} leave types", "bot_logic")
                debug_log(f"About to call session_manager.update_session...", "bot_logic")
                self.session_manager.update_session(thread_id, {'leave_types': leave_types})
                debug_log(f"Session updated successfully", "bot_logic")
            except Exception as session_error:
                debug_log(f"Error updating session: {session_error}", "general")
                import traceback
                traceback.print_exc()
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I'm having trouble managing your time-off request session. Please try again.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Check if leave type was already extracted
            try:
                debug_log(f"Checking for pre-extracted leave type...", "bot_logic")
                if 'leave_type' in extracted_data:
                    debug_log(f"Found pre-extracted leave type: {extracted_data['leave_type']}", "bot_logic")
                    # Try to map to actual leave type
                    extracted_type = extracted_data['leave_type']
                    matched_type = None
                    for lt in leave_types:
                        if extracted_type.lower() in lt.get('name', '').lower():
                            matched_type = lt
                            break
                    
                    if matched_type:
                        debug_log(f"Matched leave type: {matched_type['name']}", "bot_logic")
                        self.session_manager.update_session(thread_id, {'selected_leave_type': matched_type})
                        self._persist_timeoff_context(thread_id, session, selected_leave_type=matched_type)
                        self.session_manager.advance_session_step(thread_id, {'leave_type_confirmed': True})

                        # If both dates are already present, move straight to confirmation
                        if 'start_date' in extracted_data and 'end_date' in extracted_data:
                            start_date = self.timeoff_service.parse_date_input(extracted_data['start_date']) or extracted_data['start_date']
                            end_date = self.timeoff_service.parse_date_input(extracted_data['end_date']) or extracted_data['end_date']
                            self.session_manager.update_session(thread_id, {'start_date': start_date, 'end_date': end_date})
                            self._persist_timeoff_context(thread_id, session, start_date=start_date, end_date=end_date)
                            self.session_manager.advance_session_step(thread_id)

                            # Enforce Sick Leave mode choice before confirmation
                            try:
                                mt_name_cf = (matched_type.get('name') or '').strip()
                            except Exception:
                                mt_name_cf = ''
                            if self._mode_config_for_leave(mt_name_cf) and not self._get_leave_mode(session, mt_name_cf):
                                return self._prompt_leave_mode(thread_id, mt_name_cf)

                            def dd_slash_mm_yyyy(d: str) -> str:
                                try:
                                    return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                                except Exception:
                                    return d
                            response_text = f"Perfect! Here's your time-off request summary:\n\n"
                            response_text += f"📋 **Leave Type:** {matched_type.get('name', 'Unknown')}\n"
                            response_text += f"📅 **Start Date:** {dd_slash_mm_yyyy(start_date)}\n"
                            response_text += f"📅 **End Date:** {dd_slash_mm_yyyy(end_date)}\n"
                            response_text += f"👤 **Employee:** {employee_data.get('name', 'Unknown')}\n\n"
                            response_text += "Do you want to submit this request? Reply with 'yes' to confirm or 'no' to cancel."
                            return self._create_response(response_text, thread_id)

                        # Otherwise, first ask Sick Leave mode if applicable; else show date picker
                        try:
                            mt_name = (matched_type.get('name') or '').strip()
                        except Exception:
                            mt_name = ''
                        if self._mode_config_for_leave(mt_name) and not self._get_leave_mode(session, mt_name):
                            return self._prompt_leave_mode(thread_id, mt_name)
                        response_text = (
                            f"I'll help you request {matched_type['name']}. \n\n"
                            "You can pick dates from the calendar below or type them. Examples:\n"
                            "- 23/9 to 24/9\n"
                            "- 23/09/2025 to 24/09/2025\n"
                            "- 23-9-2025 till 24-9-2025\n"
                            "- 23rd of September till the 24th\n"
                            "- next Monday to Wednesday\n\n"
                            "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                        )
                        return self._create_response_with_datepicker(response_text, thread_id)
                else:
                    debug_log(f"No pre-extracted leave type found", "bot_logic")
            except Exception as extract_error:
                debug_log(f"Error processing extracted data: {extract_error}", "general")
                # Continue to present options
            
            # Present the main leave type options with buttons
            try:
                debug_log(f"Presenting main leave type options to user...", "bot_logic")

                # Filter to only show the core leave types we surface in the UI
                main_leave_types = []
                type_names = ['Annual Leave', 'Sick Leave', 'Unpaid Leave', 'Custom Hours']

                for type_name in type_names:
                    for lt in leave_types:
                        if lt.get('name') == type_name:
                            main_leave_types.append(lt)
                            break

                debug_log(f"Found {len(main_leave_types)} main leave types", "bot_logic")

                # Store the main leave types in session
                self.session_manager.update_session(thread_id, {'main_leave_types': main_leave_types})

                response_text = "I'll help you request time off! Please select the type of leave you need:"

                # Create response with buttons
                result = self._create_response_with_buttons(
                    response_text,
                    thread_id,
                    main_leave_types
                )
                debug_log(f"Response with buttons created successfully: {result is not None}", "bot_logic")
                return result
                
            except Exception as response_error:
                debug_log(f"Error creating leave type options response: {response_error}", "general")
                import traceback
                traceback.print_exc()
                
                # Fallback response
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I can help you request time off, but I'm having trouble loading the available leave types right now. Please contact HR for assistance.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
        except Exception as e:
            debug_log(f"Error starting time-off session: {e}", "general")
            # Clear session on error to prevent stuck state
            self.session_manager.clear_session(thread_id)
            return {
                'message': "I'd like to help with your time-off request, but I'm experiencing some technical difficulties. Please try again later.",
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 1.0,
                'model_used': self.model
            }
    
    def _continue_timeoff_session(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Continue an active time-off session"""
        try:
            # Refresh the session snapshot to pick up any updates from previous steps
            try:
                if thread_id:
                    refreshed_session = self.session_manager.get_session(thread_id)
                    if refreshed_session:
                        session = refreshed_session
            except Exception:
                pass

            # If session is already completed/cancelled, clear and do not continue
            try:
                state = session.get('state')
                if state in ['completed', 'cancelled']:
                    self.session_manager.clear_session(thread_id)
                    return None
            except Exception:
                pass
            # Check for exit commands
            message_lower = message.lower().strip()
            if message_lower in ['cancel', 'exit', 'stop', 'quit', 'nevermind', 'no thanks', 'end', 'abort', 'undo', 'no', 'n']:
                debug_log(f"User said '{message_lower}', cancelling time-off session", "bot_logic")
                try:
                    self.session_manager.cancel_session(thread_id, f"User requested to exit time-off flow with: {message_lower}")
                finally:
                    # Clear any persisted state to avoid sticky sessions
                    self._reset_timeoff_flow_state(thread_id)
                return {
                    'message': 'request cancelled, can i help you with anything else',
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }

            # Allow users to restart mid-flow with a fresh time-off request phrase
            if self._is_timeoff_start_message(message):
                debug_log("Restart phrase detected during continuation; resetting flow.", "bot_logic")
                payload = {}
                try:
                    detected = self.timeoff_service.detect_timeoff_intent(message)
                    if detected[0]:
                        payload = detected[2] or {}
                except Exception:
                    payload = {}
                self._reset_timeoff_sessions(thread_id, 'User restarted time-off flow during continuation', employee_data)
                return self._start_timeoff_session(message, thread_id, payload, employee_data)

            # Require supporting document upload (if pending) before moving on
            if self._is_supporting_doc_stage(session):
                return self._handle_supporting_document_step(message, thread_id, session, employee_data)

            # Fast-path: confirmation keywords should move straight to submission when context exists
            confirmation_tokens = {'yes', 'y', 'confirm', 'submit', 'ok', 'sure'}
            if message_lower in confirmation_tokens:
                ctx = self._resolve_timeoff_context(session)
                selected_type_confirm = ctx.get('selected_leave_type')
                start_confirm = ctx.get('start_date')
                end_confirm = ctx.get('end_date')

                if selected_type_confirm and start_confirm and end_confirm:
                    try:
                        # Ensure we are at confirmation step
                        self.session_manager.update_session(thread_id, {'step': 3})
                    except Exception:
                        pass
                    # Re-fetch once more so downstream logic sees the updated step/data
                    try:
                        if thread_id:
                            refreshed_for_confirmation = self.session_manager.get_session(thread_id)
                            if refreshed_for_confirmation:
                                session = refreshed_for_confirmation
                    except Exception:
                        pass
                    return self._handle_confirmation(message, thread_id, session, employee_data)

                # If confirmation comes too early, strictly enforce order
                if not selected_type_confirm:
                    debug_log("Confirmation received but leave type missing; prompting user to reselect.", "bot_logic")
                    try:
                        self.session_manager.update_session(thread_id, {'step': 1})
                    except Exception:
                        pass
                    leave_types = session.get('main_leave_types') or session.get('leave_types', [])
                    if leave_types:
                        prompt = "I lost track of which leave type you picked. Please choose it again:"
                        return self._create_response_with_buttons(prompt, thread_id, leave_types)
                    prompt = "I lost track of which leave type you picked. Please let me know the leave type you need."
                    return self._create_response(prompt, thread_id)

                debug_log("Confirmation received but dates missing; requesting date range again.", "bot_logic")
                try:
                    self.session_manager.update_session(thread_id, {'step': 2})
                except Exception:
                    pass
                reprompt = (
                    "I still need both the start and end date before I can submit the request. "
                    "Please send them in one message, for example '15/10/2025 to 16/10/2025'."
                )
                return self._create_response_with_datepicker(reprompt, thread_id)

            step = session.get('step', 1)
            debug_log(f"Continuing session at step {step} for thread {thread_id}", "bot_logic")

            # Enforce strict ordered steps: 1) leave type -> 2) dates -> 3) confirmation
            if step <= 1:  # Waiting for leave type selection
                return self._handle_leave_type_selection(message, thread_id, session, employee_data)
            elif step == 2:  # Waiting for date range (start and end in one message)
                # Insert leave-mode choice (Full Days vs Custom Hours) BEFORE date step when required
                try:
                    sd = session.get('data', {}) if isinstance(session, dict) else {}
                    selected_type = sd.get('selected_leave_type') or session.get('selected_leave_type') or {}
                    selected_name = (selected_type.get('name') or '').strip()
                    cfg = self._mode_config_for_leave(selected_name)
                    current_mode = self._get_leave_mode(session, selected_name) if cfg else None
                except Exception:
                    selected_name = ''
                    cfg = {}
                    current_mode = None

                # Handle response to mode-selection buttons when applicable
                try:
                    ml_raw = (message or '').strip()
                    # Normalize whitespace variants (space, NBSP, thin space) and collapse
                    ml_ws = ml_raw.replace('\u00A0', ' ').replace('\u2007', ' ').replace('\u202F', ' ')
                    import re as _re_norm
                    ml_ws = _re_norm.sub(r"\s+", " ", ml_ws).strip()
                    ml = ml_ws.lower()

                    awaiting_mode = bool(cfg) and not bool(current_mode)
                    if cfg and awaiting_mode:
                        field_key = cfg.get('session_key') or ''
                        prefix = cfg.get('button_prefix') or selected_name.upper().replace(' ', '_')
                        label = cfg.get('prompt_label') or selected_name or 'this leave type'

                        # Accept explicit widget payload format: <field_key>=VALUE
                        if field_key and ml.startswith(f"{field_key.lower()}="):
                            val = ml_raw.split('=', 1)[1].strip()
                            val_upper = val.upper()
                            if val_upper == f"{prefix}_FULL_DAYS" or val_upper == "FULL_DAYS":
                                self._store_leave_mode(thread_id, session, selected_name, 'full_days')
                                guidance = (
                                    f"Great, we will request {label} for full days. "
                                    "Please send both dates in one message, for example '15/10/2025 to 16/10/2025'."
                                )
                                return self._create_response_with_datepicker(guidance, thread_id)
                            elif val_upper == f"{prefix}_CUSTOM_HOURS" or val_upper == "CUSTOM_HOURS":
                                self._store_leave_mode(thread_id, session, selected_name, 'custom_hours')
                                guidance = (
                                    f"Great, we will request {label} for custom hours. \n\n"
                                    "Pick your date from the calendar below or type it (single day only). Examples:\n"
                                    "- 23/9\n"
                                    "- 23/09/2025\n"
                                    "- next Monday\n\n"
                                    "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                                )
                                return self._create_response_with_datepicker_single(guidance, thread_id)

                        # Normalize common variants users/typeahead may send
                        is_full_days = (
                            ml == 'full days' or ml == 'fulldays' or ml == 'full-day' or
                            bool(_re_norm.search(r"\bfull\s*[- ]?\s*days\b", ml)) or ml_raw.upper() == f"{prefix}_FULL_DAYS"
                        )
                        is_custom_hours = (
                            ml == 'custom hours' or ml == 'custom-hours' or ml == 'custom hour' or
                            bool(_re_norm.search(r"\bcustom\s*[- ]?\s*hours?\b", ml)) or ml_raw.upper() == f"{prefix}_CUSTOM_HOURS"
                        )
                        if is_full_days or is_custom_hours:
                            mode_label = 'full_days' if is_full_days else 'custom_hours'
                            debug_log(f"{label} mode selected: {'Full Days' if is_full_days else 'Custom Hours'} (raw: {ml_raw})", "bot_logic")
                            self._store_leave_mode(thread_id, session, selected_name, mode_label)
                            if mode_label == 'full_days':
                                guidance = (
                                    f"Great, we will request {label} for full days. "
                                    "Please send both dates in one message, for example '15/10/2025 to 16/10/2025'."
                                )
                                return self._create_response_with_datepicker(guidance, thread_id)
                            else:
                                guidance = (
                                    f"Great, we will request {label} for custom hours. \n\n"
                                    "Pick your date from the calendar below or type it (single day only). Examples:\n"
                                    "- 23/9\n"
                                    "- 23/09/2025\n"
                                    "- next Monday\n\n"
                                    "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                                )
                                return self._create_response_with_datepicker_single(guidance, thread_id)
                except Exception as e:
                    debug_log(f"Error handling leave mode selection: {e}", "general")
                    import traceback
                    traceback.print_exc()

                # If a mode is required but not chosen yet, prompt the user
                if cfg and not current_mode and selected_name:
                    return self._prompt_leave_mode(thread_id, selected_name)

                # Safety: if no leave type was selected (e.g., residue from a previous flow), reset to step 1
                try:
                    sd = session.get('data', {}) if isinstance(session, dict) else {}
                    selected_type = sd.get('selected_leave_type') or session.get('selected_leave_type')
                except Exception:
                    selected_type = None
                if not selected_type:
                    debug_log("Step=2 but no selected_leave_type found; resetting to step 1", "bot_logic")
                    self.session_manager.update_session(thread_id, {'step': 1})
                    return self._handle_leave_type_selection(message, thread_id, session, employee_data)
                return self._handle_date_range_input(message, thread_id, session, employee_data)
            elif step >= 3:  # Waiting for confirmation
                return self._handle_confirmation(message, thread_id, session, employee_data)
            else:
                debug_log(f"Session in unknown state (step {step}), restarting", "bot_logic")
                # Session in unknown state, restart
                self.session_manager.clear_session(thread_id)
                return None

        except Exception as e:
            debug_log(f"Error continuing time-off session: {e}", "general")
            import traceback
            traceback.print_exc()
            # Clear the broken session
            self.session_manager.clear_session(thread_id)
            return None
    
    def _handle_leave_type_selection(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle leave type selection step"""
        if self._is_timeoff_start_message(message):
            debug_log("Restart phrase detected during leave type selection; restarting flow.", "bot_logic")
            return self._restart_timeoff_flow(message, thread_id, employee_data, 'User restarted during leave type selection')

        session_data = session.get('data', {})
        # Check for main leave types first (the 3 button options)
        leave_types = session_data.get('main_leave_types', [])
        if not leave_types:
            # Fallback to full leave types
            leave_types = session_data.get('leave_types', [])
        if not leave_types:
            # Also check at session root level
            leave_types = session.get('leave_types', []) or session.get('main_leave_types', [])

        debug_log(f"Found {len(leave_types)} leave types in session", "bot_logic")
        
        # Validate that we still have leave types
        if not leave_types:
            debug_log("No leave types found in session data; forcing a fresh start.", "bot_logic")
            self._reset_timeoff_sessions(thread_id, 'Missing leave types', employee_data)
            restarted = self._start_timeoff_session("time off", thread_id, {}, employee_data)
            if restarted:
                return restarted
            return {
                'message': "I'm sorry, but I'm having trouble accessing the leave types. Please try requesting time off again or contact HR for assistance.",
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 1.0,
                'model_used': self.model
            }
        
        message_clean = message.strip()
        
        # Helper: primary types from session or filter
        def _primary_types(all_types: list) -> list:
            try:
                names = ['Annual Leave', 'Sick Leave', 'Unpaid Leave', 'Custom Hours']
                result = []
                for nm in names:
                    for lt in all_types:
                        if lt.get('name') == nm:
                            result.append(lt)
                            break
                return result or all_types
            except Exception:
                return all_types

        # Check if user provided a number
        try:
            choice_num = int(message_clean)
            if 1 <= choice_num <= len(leave_types):
                selected_type = leave_types[choice_num - 1]
                # If dates already exist (e.g., after reselect on confirmation), jump directly to confirmation
                ctx = self._resolve_timeoff_context(session)
                existing_start = ctx.get('start_date')
                existing_end = ctx.get('end_date')
                if existing_start and existing_end:
                    # If Sick Leave, enforce mode prompt before confirmation
                    try:
                        sel_name = (selected_type.get('name') or '').strip()
                    except Exception:
                        sel_name = ''
                    if self._mode_config_for_leave(sel_name) and not self._get_leave_mode(session, sel_name):
                        try:
                            self.session_manager.update_session(
                                thread_id,
                                {
                                    'selected_leave_type': selected_type,
                                    'step': 2
                                }
                            )
                        except Exception:
                            pass
                        self._persist_timeoff_context(thread_id, session, selected_leave_type=selected_type)
                        return self._prompt_leave_mode(thread_id, sel_name)
                    try:
                        self.session_manager.update_session(
                            thread_id,
                            {
                                'selected_leave_type': selected_type,
                                'start_date': existing_start,
                                'end_date': existing_end,
                                'step': 3,
                                # clear pendings
                                'data': {**session.get('data', {}), 'pending_start_date': None, 'pending_end_date': None}
                            }
                        )
                    except Exception:
                        pass
                    self._persist_timeoff_context(
                        thread_id,
                        session,
                        selected_leave_type=selected_type,
                        start_date=existing_start,
                        end_date=existing_end
                    )
                    def _fmt(d: str) -> str:
                        try:
                            return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                        except Exception:
                            return d
                    
                    # Get remaining leave time
                    remaining_leave_text = ""
                    try:
                        if self.leave_balance_service and employee_data and employee_data.get('id'):
                            employee_id = employee_data.get('id')
                            leave_type_name = selected_type.get('name', 'Unknown')
                            remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, leave_type_name)
                            if not error:
                                # Always show balance, even if 0 (for specific leave type)
                                if remaining:
                                    formatted = self.leave_balance_service.format_remaining_leave_message(remaining)
                                    if formatted:
                                        # Format as "**Available Annual Leave:** X days" (only label bold, not days)
                                        import re
                                        formatted_bold_label = re.sub(r'(Available [^:]+):', r'**\1:**', formatted)
                                        remaining_leave_text = f"\n⏰ {formatted_bold_label}"
                                elif leave_type_name:
                                    # Show 0 days for specific leave type with no allocations
                                    remaining_leave_text = f"\n⏰ **Available {leave_type_name}:** 0 days"
                    except Exception as e:
                        debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                        # Continue without remaining leave info if there's an error
                    
                    response_text = (
                        "Perfect! Here's your time-off request summary:\n\n"
                        f"📋 **Leave Type:** {selected_type.get('name', 'Unknown')}\n"
                        f"📅 **Start Date:** {_fmt(existing_start)}\n"
                        f"📅 **End Date:** {_fmt(existing_end)}\n"
                        f"👤 **Employee:** {(employee_data or {}).get('name', 'Unknown')}{remaining_leave_text}\n"
                        "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    )
                    buttons = [
                        {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                        {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                    ]
                    return self._create_response_with_choice_buttons(response_text, thread_id, buttons)
                # Otherwise, proceed to next step (for Sick Leave: ask mode first)
                self.session_manager.update_session(thread_id, {'selected_leave_type': selected_type})
                self._persist_timeoff_context(thread_id, session, selected_leave_type=selected_type)
                self.session_manager.advance_session_step(thread_id)
                try:
                    sel_name2 = (selected_type.get('name') or '').strip()
                except Exception:
                    sel_name2 = ''
                if self._mode_config_for_leave(sel_name2) and not self._get_leave_mode(session, sel_name2):
                    return self._prompt_leave_mode(thread_id, sel_name2)
                response_text = f"Great! You've selected {selected_type['name']}. \n\nYou can pick dates from the calendar below or type them. Examples:\n- 23/9 to 24/9\n- 23/09/2025 to 24/09/2025\n- 23-9-2025 till 24-9-2025\n- 23rd of September till the 24th\n- next Monday to Wednesday\n\nDefaults: I assume DD/MM, current month and year unless you specify otherwise."
                return self._create_response_with_datepicker(response_text, thread_id)
        except ValueError:
            pass
        
        # Try to match by name (improved matching)
        message_lower = message_clean.lower()
        best_match = None
        best_score = 0
        
        for leave_type in leave_types:
            leave_name = leave_type.get('name', '').lower()
            # Check for exact word matches
            name_words = leave_name.split()
            message_words = message_lower.split()
            
            matches = 0
            for word in name_words:
                if word in message_words:
                    matches += 1
            
            # Also check for partial name inclusion
            if leave_name in message_lower or message_lower in leave_name:
                matches += len(name_words)
            
            if matches > best_score:
                best_score = matches
                best_match = leave_type
        
        # If we found a good match, use it
        if best_match and best_score > 0:
            # If dates already exist (e.g., after reselect on confirmation), jump directly to confirmation
            ctx = self._resolve_timeoff_context(session)
            existing_start = ctx.get('start_date')
            existing_end = ctx.get('end_date')
            if existing_start and existing_end:
                # If Sick Leave, enforce mode prompt before confirmation
                try:
                    bm_name = (best_match.get('name') or '').strip()
                except Exception:
                    bm_name = ''
                if self._mode_config_for_leave(bm_name) and not self._get_leave_mode(session, bm_name):
                    try:
                        self.session_manager.update_session(
                            thread_id,
                            {
                                'selected_leave_type': best_match,
                                'step': 2
                            }
                        )
                    except Exception:
                        pass
                    self._persist_timeoff_context(thread_id, session, selected_leave_type=best_match)
                    return self._prompt_leave_mode(thread_id, bm_name)
                try:
                    self.session_manager.update_session(
                        thread_id,
                        {
                            'selected_leave_type': best_match,
                            'start_date': existing_start,
                            'end_date': existing_end,
                            'step': 3,
                            'data': {**session.get('data', {}), 'pending_start_date': None, 'pending_end_date': None}
                        }
                    )
                except Exception:
                    pass
                self._persist_timeoff_context(
                    thread_id,
                    session,
                    selected_leave_type=best_match,
                    start_date=existing_start,
                    end_date=existing_end
                )
                def _fmt(d: str) -> str:
                    try:
                        return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                    except Exception:
                        return d
                
                # Get remaining leave time
                remaining_leave_text = ""
                try:
                    if self.leave_balance_service and employee_data and employee_data.get('id'):
                        employee_id = employee_data.get('id')
                        leave_type_name = best_match.get('name', 'Unknown')
                        remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, leave_type_name)
                        if not error:
                            # Always show balance, even if 0 (for specific leave type)
                            if remaining:
                                remaining_leave_text = self.leave_balance_service.format_remaining_leave_message(remaining)
                                if remaining_leave_text:
                                    remaining_leave_text = f"\n⏰ **{remaining_leave_text}**\n"
                            elif leave_type_name:
                                # Show 0 days for specific leave type with no allocations
                                remaining_leave_text = f"\n⏰ **Available {leave_type_name}:** 0 days\n"
                except Exception as e:
                    debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                    # Continue without remaining leave info if there's an error
                
                response_text = (
                    "Perfect! Here's your time-off request summary:\n\n"
                    f"📋 **Leave Type:** {best_match.get('name', 'Unknown')}\n"
                    f"📅 **Start Date:** {_fmt(existing_start)}\n"
                    f"📅 **End Date:** {_fmt(existing_end)}\n"
                    f"👤 **Employee:** {(employee_data or {}).get('name', 'Unknown')}{remaining_leave_text}\n\n"
                    "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                )
                buttons = [
                    {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                    {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                ]
                return self._create_response_with_choice_buttons(response_text, thread_id, buttons)

            # Otherwise, proceed to next step (for Sick Leave: ask mode first)
            self.session_manager.update_session(thread_id, {'selected_leave_type': best_match})
            self._persist_timeoff_context(thread_id, session, selected_leave_type=best_match)
            self.session_manager.advance_session_step(thread_id)
            
            # If Half Days, show single-date picker; else show range picker
            try:
                is_halfday = False
                if self.halfday_service and isinstance(best_match, dict):
                    is_halfday = self.halfday_service.is_halfday(best_match) or (best_match.get('name') == getattr(self.halfday_service, 'HALF_DAY_NAME', 'Custom Hours'))
            except Exception:
                is_halfday = False

            if is_halfday:
                response_text = (
                    f"Perfect! You've selected {best_match['name']}. \n\n"
                    "Pick your date from the calendar below or type it (single day only). Examples:\n"
                    "- 23/9\n"
                    "- 23/09/2025\n"
                    "- next Monday\n\n"
                    "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                )
                return self._create_response_with_datepicker_single(response_text, thread_id)
            else:
                try:
                    bm_name2 = (best_match.get('name') or '').strip()
                except Exception:
                    bm_name2 = ''
                if self._mode_config_for_leave(bm_name2) and not self._get_leave_mode(session, bm_name2):
                    return self._prompt_leave_mode(thread_id, bm_name2)
                response_text = (
                    f"Perfect! You've selected {best_match['name']}. \n\nYou can pick dates from the calendar below or type them. Examples:\n"
                    "- 23/9 to 24/9\n"
                    "- 23/09/2025 to 24/09/2025\n"
                    "- 23-9-2025 till 24-9-2025\n"
                    "- 23rd of September till the 24th\n"
                    "- next Monday to Wednesday\n\n"
                    "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                )
                return self._create_response_with_datepicker(response_text, thread_id)
        
        # No match found - before falling back, check if the user already provided dates.
        try:
            # If the message looks like a date range or a single date, treat it as step 2 input.
            dr = self.timeoff_service.parse_date_range(message)
            if dr:
                return self._handle_date_range_input(message, thread_id, session, employee_data)
            single = self.timeoff_service.parse_date_input(message)
            if single:
                # Treat as same-day range
                return self._handle_date_range_input(message, thread_id, session, employee_data)
        except Exception:
            pass

        # Still no match - before guidance, if we already captured dates, move to confirmation path by asking for a type explicitly with buttons
        try:
            sd3 = session.get('data', {}) if isinstance(session, dict) else {}
            ctx3 = self._resolve_timeoff_context(session)
            start3 = ctx3.get('start_date') or sd3.get('pending_start_date')
            end3 = ctx3.get('end_date') or sd3.get('pending_end_date')
            if start3 and end3:
                prim = _primary_types(leave_types)
                msg3 = "I have your dates. Please choose a leave type to continue:"
                return self._create_response_with_buttons(msg3, thread_id, prim)
        except Exception:
            pass

        # Still no match - provide helpful guidance
        # Prefer showing the primary options again instead of dumping the full list
        try:
            primary_names = ['Annual Leave', 'Sick Leave', 'Unpaid Leave', 'Custom Hours']
            primary = []
            for nm in primary_names:
                for lt in leave_types:
                    if lt.get('name') == nm:
                        primary.append(lt)
                        break
            if primary:
                msg = "I didn't quite catch that. Please choose one of the leave types below:"
                return self._create_response_with_buttons(msg, thread_id, primary)
        except Exception:
            pass
        # Fallback to formatted list only if we cannot find the primary set
        leave_types_text = self.timeoff_service.format_leave_types_for_user(leave_types)
        response_text = f"I didn't quite catch that. Please select from the available options:\n\n{leave_types_text}\nYou can type the number (1, 2, 3, etc.) or the name of the leave type (like 'annual' or 'sick')."
        return self._create_response(response_text, thread_id)
    
    def _handle_start_date_input(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle start date input step"""
        if self._is_timeoff_start_message(message):
            debug_log("Restart phrase detected during start date input; restarting flow.", "bot_logic")
            return self._restart_timeoff_flow(message, thread_id, employee_data, 'User restarted during start date step')

        parsed_date = self.timeoff_service.parse_date_input(message)
        
        if parsed_date:
            self.session_manager.update_session(thread_id, {'start_date': parsed_date})
            self._persist_timeoff_context(thread_id, session, start_date=parsed_date)
            self.session_manager.advance_session_step(thread_id)
            
            response_text = f"Got it! Start date: {parsed_date}\n\nWhat's your end date? Please use DD/MM/YYYY format."
            return self._create_response(response_text, thread_id)
        else:
            response_text = "I couldn't understand that date format. Please provide the start date in DD/MM/YYYY format (e.g., 25/12/2024)."
            return self._create_response(response_text, thread_id)
    
    def _handle_end_date_input(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle end date input step"""
        if self._is_timeoff_start_message(message):
            debug_log("Restart phrase detected during end date input; restarting flow.", "bot_logic")
            return self._restart_timeoff_flow(message, thread_id, employee_data, 'User restarted during end date step')

        parsed_date = self.timeoff_service.parse_date_input(message)
        
        if parsed_date:
            session_data = session.get('data', {})
            context_data = {}
            if isinstance(session_data, dict):
                context_data = session_data.get('timeoff_context', {}) or {}
            # Check both locations for start_date and selected_leave_type
            start_date = (
                session_data.get('start_date')
                or session.get('start_date')
                or context_data.get('start_date')
            )
            selected_type = (
                session_data.get('selected_leave_type')
                or session.get('selected_leave_type', {})
                or context_data.get('selected_leave_type')
            )
            # Robustly resolve employee data from argument or session
            resolved_employee = (
                employee_data
                or session_data.get('employee_data')
                or context_data.get('employee_data')
                or session.get('employee_data')
                or {}
            )
            
            debug_log(f"End date processing - start_date: {start_date}, selected_type: {selected_type.get('name', 'None') if selected_type else 'None'}", "bot_logic")
            
            # Validate end date is after start date
            if start_date and parsed_date < start_date:
                response_text = f"The end date ({parsed_date}) cannot be before the start date ({start_date}). Please provide a valid end date."
                return self._create_response(response_text, thread_id)
            
            self.session_manager.update_session(thread_id, {'end_date': parsed_date})
            self._persist_timeoff_context(
                thread_id,
                session,
                start_date=start_date,
                end_date=parsed_date,
                selected_leave_type=selected_type
            )
            self.session_manager.advance_session_step(thread_id)
            
            # Show confirmation (display dates as DD/MM/YYYY)
            def dd_mm_yyyy(d: str) -> str:
                try:
                    return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                except Exception:
                    return d
            
            # Get remaining leave time
            remaining_leave_text = ""
            try:
                if self.leave_balance_service and resolved_employee and resolved_employee.get('id'):
                    employee_id = resolved_employee.get('id')
                    leave_type_name = selected_type.get('name', 'Unknown') if selected_type else 'Unknown'
                    remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, leave_type_name)
                    if not error and remaining:
                        remaining_leave_text = self.leave_balance_service.format_remaining_leave_message(remaining)
                        if remaining_leave_text:
                            remaining_leave_text = f"\n⏰ **{remaining_leave_text}**\n"
            except Exception as e:
                debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                # Continue without remaining leave info if there's an error
            
            response_text = f"Perfect! Here's your time-off request summary:\n\n"
            response_text += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if selected_type else 'Unknown'}\n"
            response_text += f"📅 **Start Date:** {dd_mm_yyyy(start_date) if start_date else 'Unknown'}\n"
            response_text += f"📅 **End Date:** {dd_mm_yyyy(parsed_date)}\n"
            response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}{remaining_leave_text}\n\n"
            response_text += "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
            buttons = [
                {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
            ]
            debug_log(f"Returning confirmation response with buttons: {response_text[:100]}...", "bot_logic")
            return self._create_response_with_choice_buttons(response_text, thread_id, buttons)
        else:
            response_text = "I couldn't understand that date format. Please provide the end date in DD/MM/YYYY format (e.g., 27/12/2024)."
            return self._create_response(response_text, thread_id)

    def _handle_date_range_input(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle combined date range input (start and end in one message)"""
        if self._is_timeoff_start_message(message):
            debug_log("Restart phrase detected during date range input; restarting flow.", "bot_logic")
            return self._restart_timeoff_flow(message, thread_id, employee_data, 'User restarted during date range step')

        # Refresh session snapshot to capture any updates from previous steps
        try:
            if thread_id:
                refreshed_session = self.session_manager.get_session(thread_id)
                if refreshed_session:
                    session = refreshed_session
        except Exception:
            pass

        # First, if the user typed confirmation terms here, jump to confirmation when context is ready
        try:
            ml = (message or '').strip().lower()
            if ml in {'yes','y','confirm','submit','ok','sure','no','n','cancel','abort','stop','exit','quit'}:
                sd_cf = session.get('data', {}) if isinstance(session, dict) else {}
                sel_type_cf = sd_cf.get('selected_leave_type') or session.get('selected_leave_type')
                start_cf = sd_cf.get('start_date') or session.get('start_date')
                end_cf = sd_cf.get('end_date') or session.get('end_date')
                if sel_type_cf and start_cf and end_cf:
                    try:
                        self.session_manager.update_session(thread_id, {'step': 3})
                    except Exception:
                        pass
                    return self._handle_confirmation(message, thread_id, session, employee_data)
        except Exception:
            pass

        # If the message has no date hints at all, just re-show the date picker without an error
        try:
            s_msg = (message or '').strip().lower()
            import re as _re
            has_digits = bool(_re.search(r"\d", s_msg))
            has_connectors = any(t in s_msg for t in [' to ', ' until ', ' till ', '-', '–', '—'])
            has_words = any(t in s_msg for t in ['tomorrow', 'today', 'next ', 'monday','tuesday','wednesday','thursday','friday','saturday','sunday'])
            has_widget_key = 'timeoff_date_range=' in s_msg or 'embassy_date_range=' in s_msg or 'overtime_date_range=' in s_msg or 'reimbursement_expense_date=' in s_msg
            if not (has_digits or has_connectors or has_words or has_widget_key):
                guidance = (
                    "Please send both dates in one message, for example '15/10/2025 to 16/10/2025'."
                )
                return self._create_response_with_datepicker(guidance, thread_id)
        except Exception:
            pass

        # Detect if current flow is Half Days or Custom Hours (Sick/Unpaid) to enforce single-day constraint
        session_data_for_type = session.get('data', {})
        selected_type_for_validation = session_data_for_type.get('selected_leave_type') or session.get('selected_leave_type', {})
        
        # Check if this is a custom-hours mode for the current leave type
        selected_name_check = ''
        is_custom_hours_mode = False
        try:
            selected_name_check = (selected_type_for_validation.get('name') or '').strip()
            is_custom_hours_mode = self._is_custom_hours_mode(session, selected_name_check)
        except Exception:
            is_custom_hours_mode = False
        custom_label = selected_name_check or 'This leave'
        
        is_halfday_flow = False
        try:
            if self.halfday_service and isinstance(selected_type_for_validation, dict):
                is_halfday_flow = self.halfday_service.is_halfday(selected_type_for_validation) or (
                    selected_type_for_validation.get('name') == getattr(self.halfday_service, 'HALF_DAY_NAME', 'Custom Hours')
                )
        except Exception:
            is_halfday_flow = False
        
        # Treat Sick/Unpaid Custom Hours like halfday flow (single date + hours)
        is_single_date_flow = is_halfday_flow or is_custom_hours_mode

        # Accept widget format strictly first: "DD/MM/YYYY to DD/MM/YYYY"
        try:
            raw = (message or '').strip()
            # Accept multiple widget keys
            for key in ['reimbursement_expense_date=', 'embassy_date_range=', 'overtime_date_range=', 'timeoff_date_range=']:
                if raw.lower().startswith(key):
                    raw = raw.split('=', 1)[1].strip()
                    if not raw:
                        # Empty widget payload: re-open picker without error
                        return self._create_response_with_datepicker(
                            "Please select a start and end date.", thread_id
                        )
                    break
            if ' to ' in raw and len(raw.split(' to ')) == 2:
                a, b = [p.strip() for p in raw.split(' to ')]
                from datetime import datetime as _dt
                for dt in (a, b):
                    _dt.strptime(dt, '%d/%m/%Y')
                # widget-provided format is valid; trust it and store
                start_date = _dt.strptime(a, '%d/%m/%Y').strftime('%Y-%m-%d')
                end_date = _dt.strptime(b, '%d/%m/%Y').strftime('%Y-%m-%d')
                if is_single_date_flow and start_date != end_date:
                    try:
                        self.session_manager.update_session(thread_id, {'step': 2})
                    except Exception:
                        pass
                    msg = (
                        "Custom Hours are limited to one day. Please pick a single date."
                        if is_halfday_flow
                        else f"{custom_label} Custom Hours are limited to one day. Please pick a single date."
                    )
                    return self._create_response_with_datepicker_single(msg, thread_id)
                return self._process_captured_dates(
                    thread_id,
                    session,
                    employee_data,
                    start_date,
                    end_date,
                    is_single_date_flow
                )
        except Exception:
            pass

        result = self.timeoff_service.parse_date_range(message)
        if result:
            start_date, end_date = result
            # If single-date flow (Half Day or Sick Custom Hours), only allow single day
            if is_single_date_flow and start_date != end_date:
                msg = (
                    "For Half Days, you can only select one day. Please pick a single date."
                    if is_halfday_flow
                    else f"For {custom_label} Custom Hours, you can only select one day. Please pick a single date."
                )
                return self._create_response_with_datepicker_single(msg, thread_id)
            # Validate chronological order (already ensured) and store
            return self._process_captured_dates(
                thread_id,
                session,
                employee_data,
                start_date,
                end_date,
                is_single_date_flow
            )
        else:
            # Holistic fallback: accept a single date and treat it as a same-day range
            single = self.timeoff_service.parse_date_input(message)
            if single:
                return self._process_captured_dates(
                    thread_id,
                    session,
                    employee_data,
                    single,
                    single,
                    is_single_date_flow
                )

            response_text = (
                "I couldn't parse the date range. Please send both dates in one message. Examples:\n"
                "- 23/9 to 24/9\n"
                "- 23/09/2025 to 24/09/2025\n"
                "- 23-9-2025 till 24-9-2025\n"
                "- 23rd of September till the 24th\n"
                "- next Monday to Wednesday\n\n"
                "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
            )
            return self._create_response(response_text, thread_id)
    
    def _handle_confirmation(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle final confirmation step"""
        # Ensure we operate on the freshest session data (the handler might receive a stale snapshot)
        try:
            if thread_id:
                refreshed_session = self.session_manager.get_session(thread_id)
                if refreshed_session:
                    session = refreshed_session
        except Exception:
            pass

        message_lower = message.lower().strip()
        debug_log(f"Handling confirmation - message: '{message_lower}'", "bot_logic")

        # Confirm submit: ignore step number; rely on resolved context instead
        if message_lower in ['yes', 'y', 'confirm', 'submit', 'ok', 'sure']:
            try:
                ctx_now = self._resolve_timeoff_context(session)
            except Exception:
                ctx_now = {'selected_leave_type': None, 'start_date': None, 'end_date': None}

            sel_type_now = (ctx_now or {}).get('selected_leave_type')
            start_now = (ctx_now or {}).get('start_date')
            end_now = (ctx_now or {}).get('end_date')

            if sel_type_now and start_now and end_now:
                try:
                    self.session_manager.update_session(thread_id, {'step': 3})
                except Exception:
                    pass
                debug_log(f"User confirmed submission with full context present; submitting.", "bot_logic")
                return self._submit_timeoff_request(thread_id, session, employee_data, self.get_current_odoo_session())

            # Missing leave type
            if not sel_type_now:
                try:
                    self.session_manager.update_session(thread_id, {'step': 1})
                except Exception:
                    pass
                leave_types = session.get('main_leave_types') or session.get('leave_types', [])
                if leave_types:
                    prompt = "I lost track of which leave type you picked. Please choose it again:"
                    return self._create_response_with_buttons(prompt, thread_id, leave_types)
                return self._create_response("I lost track of which leave type you picked. Please tell me the leave type.", thread_id)

            # Missing dates
            try:
                self.session_manager.update_session(thread_id, {'step': 2})
            except Exception:
                pass
            reprompt = (
                "I still need both the start and end date before I can submit the request. "
                "Please send them in one message, for example '15/10/2025 to 16/10/2025'."
            )
            return self._create_response_with_datepicker(reprompt, thread_id)
        elif message_lower in ['no', 'n', 'cancel', 'abort', 'stop', 'exit', 'quit', 'nevermind', 'end', 'undo']:
            # Cancel the request with standardized message
            debug_log(f"User cancelled submission", "bot_logic")
            try:
                self.session_manager.cancel_session(thread_id, 'User cancelled at confirmation')
            finally:
                self._reset_timeoff_flow_state(thread_id)
            response_text = 'request cancelled, can i help you with anything else'
            return self._create_response(response_text, thread_id)
        else:
            # Allow updating dates directly at confirmation: parse date range or single date
            try:
                # 1) Date range like "DD/MM/YYYY to DD/MM/YYYY"
                dr_conf = self.timeoff_service.parse_date_range(message)
                if dr_conf:
                    start_update, end_update = dr_conf
                    self.session_manager.update_session(thread_id, {'start_date': start_update, 'end_date': end_update})
                    # Persist in context
                    try:
                        ctx_now = self._resolve_timeoff_context(session)
                        sel_type_now = (ctx_now or {}).get('selected_leave_type') or {}
                    except Exception:
                        sel_type_now = {}

                    self._persist_timeoff_context(
                        thread_id,
                        session,
                        selected_leave_type=sel_type_now,
                        start_date=start_update,
                        end_date=end_update
                    )

                    def dd_mm_yyyy(d: str) -> str:
                        try:
                            return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                        except Exception:
                            return d

                    # Build confirmation summary with updated dates
                    session_data = session.get('data', {})
                    context_data = session_data.get('timeoff_context', {}) if isinstance(session_data, dict) else {}
                    resolved_employee = (
                        employee_data
                        or session_data.get('employee_data')
                        or context_data.get('employee_data')
                        or session.get('employee_data')
                        or {}
                    )
                    lt = sel_type_now.get('name', 'Unknown') if isinstance(sel_type_now, dict) else 'Unknown'
                    
                    # Get remaining leave time (skip for unpaid leave - unlimited)
                    remaining_leave_text = ""
                    # Check if this is unpaid leave (unlimited, no balance to show)
                    is_unpaid_leave_check = (
                        isinstance(sel_type_now, dict) and 
                        sel_type_now.get('name') == 'Unpaid Leave'
                    )
                    
                    if not is_unpaid_leave_check:
                        try:
                            if self.leave_balance_service and resolved_employee and resolved_employee.get('id'):
                                employee_id = resolved_employee.get('id')
                                remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, lt)
                                if not error and remaining:
                                    remaining_leave_text = self.leave_balance_service.format_remaining_leave_message(remaining)
                                    if remaining_leave_text:
                                        remaining_leave_text = f"\n⏰ **{remaining_leave_text}**\n"
                        except Exception as e:
                            debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                            # Continue without remaining leave info if there's an error
                    
                    response_text = (
                        "Great, noted your dates. Here's your time-off request summary:\n\n"
                        f"📋 **Leave Type:** {lt}\n"
                        f"📅 **Start Date:** {dd_mm_yyyy(start_update)}\n"
                        f"📅 **End Date:** {dd_mm_yyyy(end_update)}\n"
                        f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}{remaining_leave_text}\n\n"
                        "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    )
                    buttons = [
                        {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                        {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                    ]
                    return self._create_response_with_choice_buttons(response_text, thread_id, buttons)

                # 2) Single date: treat as same-day range
                single_conf = self.timeoff_service.parse_date_input(message)
                if single_conf:
                    self.session_manager.update_session(thread_id, {'start_date': single_conf, 'end_date': single_conf})
                    try:
                        ctx_now = self._resolve_timeoff_context(session)
                        sel_type_now = (ctx_now or {}).get('selected_leave_type') or {}
                    except Exception:
                        sel_type_now = {}

                    self._persist_timeoff_context(
                        thread_id,
                        session,
                        selected_leave_type=sel_type_now,
                        start_date=single_conf,
                        end_date=single_conf
                    )

                    def _dd_mm_yyyy(d: str) -> str:
                        try:
                            return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                        except Exception:
                            return d
                    session_data = session.get('data', {})
                    context_data = session_data.get('timeoff_context', {}) if isinstance(session_data, dict) else {}
                    resolved_employee = (
                        employee_data
                        or session_data.get('employee_data')
                        or context_data.get('employee_data')
                        or session.get('employee_data')
                        or {}
                    )
                    lt = sel_type_now.get('name', 'Unknown') if isinstance(sel_type_now, dict) else 'Unknown'
                    
                    # Get remaining leave time (skip for unpaid leave - unlimited)
                    remaining_leave_text = ""
                    # Check if this is unpaid leave (unlimited, no balance to show)
                    is_unpaid_leave_check = (
                        isinstance(sel_type_now, dict) and 
                        sel_type_now.get('name') == 'Unpaid Leave'
                    )
                    
                    if not is_unpaid_leave_check:
                        try:
                            if self.leave_balance_service and resolved_employee and resolved_employee.get('id'):
                                employee_id = resolved_employee.get('id')
                                remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, lt)
                                if not error and remaining:
                                    remaining_leave_text = self.leave_balance_service.format_remaining_leave_message(remaining)
                                    if remaining_leave_text:
                                        remaining_leave_text = f"\n⏰ **{remaining_leave_text}**\n"
                        except Exception as e:
                            debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                            # Continue without remaining leave info if there's an error
                    
                    response_text = (
                        "Great, noted your date. Here's your time-off request summary:\n\n"
                        f"📋 **Leave Type:** {lt}\n"
                        f"📅 **Start Date:** { _dd_mm_yyyy(single_conf)}\n"
                        f"📅 **End Date:** { _dd_mm_yyyy(single_conf)}\n"
                        f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}{remaining_leave_text}\n\n"
                        "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    )
                    buttons = [
                        {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                        {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                    ]
                    return self._create_response_with_choice_buttons(response_text, thread_id, buttons)

                # 3) If we receive hour_from/hour_to structured message OR a natural hour range, capture and show summary
                # (existing behavior below)
                
                parsed_from, parsed_to = self._parse_hour_range_text(message)
                has_structured = ('hour_from=' in message and 'hour_to=' in message)
                if has_structured or (parsed_from and parsed_to):
                    def _get_param(k: str, s: str) -> str:
                        try:
                            parts = {p.split('=')[0]: p.split('=')[1] for p in s.split('&') if '=' in p}
                            return parts.get(k, '')
                        except Exception:
                            return ''
                    raw_from = _get_param('hour_from', message) if has_structured else parsed_from
                    raw_to = _get_param('hour_to', message) if has_structured else parsed_to

                    session_data = session.get('data', {})
                    context_data = {}
                    if isinstance(session_data, dict):
                        context_data = session_data.get('timeoff_context', {}) or {}

                    self.session_manager.update_session(thread_id, {'hour_from': raw_from, 'hour_to': raw_to})
                    self._persist_timeoff_context(thread_id, session, hour_from=raw_from, hour_to=raw_to)

                    selected_type = (
                        session_data.get('selected_leave_type')
                        or session.get('selected_leave_type', {})
                        or context_data.get('selected_leave_type')
                        or {}
                    )
                    start_date = (
                        session_data.get('start_date')
                        or session.get('start_date')
                        or context_data.get('start_date')
                    )
                    end_date = (
                        session_data.get('end_date')
                        or session.get('end_date')
                        or context_data.get('end_date')
                    )
                    resolved_employee = (
                        employee_data
                        or session_data.get('employee_data')
                        or context_data.get('employee_data')
                        or session.get('employee_data')
                        or {}
                    )

                    def dd_mm_yyyy(d: str) -> str:
                        try:
                            return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                        except Exception:
                            return d

                    # Half Day policy: max 4 hours
                    try:
                        hf = self._hour_key_to_float(raw_from)
                        ht = self._hour_key_to_float(raw_to)
                        if hf == hf and ht == ht and (ht - hf) > 4.0 - 1e-9:
                            policy_text = (
                                "Per Prezlab policy, the maximum Half Day leave duration is 4 hours. "
                                "Please enter hours less than or equal to 4."
                            )
                            return self._create_response_with_hour_picker(policy_text, thread_id)
                    except Exception:
                        pass

                    # Determine display name for leave type and base name for balance calculation
                    base_leave_type_name = selected_type.get('name', 'Custom Hours')
                    leave_type_name = base_leave_type_name
                    
                    # Check if this is Half Days (Custom Hours) flow - should use Annual Leave for balance
                    is_halfday_flow = False
                    try:
                        if self.halfday_service and isinstance(selected_type, dict):
                            is_halfday_flow = self.halfday_service.is_halfday(selected_type)
                            if is_halfday_flow:
                                # For Half Days, use Annual Leave for balance calculation
                                base_leave_type_name = 'Annual Leave'
                    except Exception:
                        pass
                    
                    try:
                        mode_for_name = self._get_leave_mode(session, base_leave_type_name)
                        if mode_for_name == 'custom_hours':
                            leave_type_name = f"{base_leave_type_name} (Custom Hours)"
                    except Exception:
                        pass
                    
                    updated_session = session
                    try:
                        latest = self.session_manager.get_session(thread_id)
                        if latest:
                            updated_session = latest
                    except Exception:
                        pass

                    if self._requires_supporting_document(updated_session) and not self._has_supporting_document(updated_session):
                        try:
                            data_payload = dict(updated_session.get('data', {}) or {})
                            ctx_payload = dict((data_payload.get('timeoff_context') or {}))
                            ctx_payload.update({
                                'supporting_doc_required': True,
                                'supporting_doc_uploaded': False
                            })
                            data_payload['timeoff_context'] = ctx_payload
                            data_payload['supporting_doc_required'] = True
                            data_payload['supporting_doc_uploaded'] = False
                            self.session_manager.update_session(thread_id, {
                                'supporting_doc_required': True,
                                'supporting_doc_uploaded': False,
                                'data': data_payload
                            })
                        except Exception:
                            pass
                        return self._prompt_supporting_document_upload(thread_id)
                    
                    # Get remaining leave time (skip for unpaid leave - unlimited)
                    remaining_leave_text = ""
                    # Check if this is unpaid leave (unlimited, no balance to show)
                    is_unpaid_leave_check = (
                        isinstance(selected_type, dict) and 
                        selected_type.get('name') == 'Unpaid Leave'
                    )
                    
                    if not is_unpaid_leave_check:
                        try:
                            if self.leave_balance_service and resolved_employee and resolved_employee.get('id'):
                                employee_id = resolved_employee.get('id')
                                # Use base_leave_type_name for balance calculation (e.g., "Annual Leave" for Half Days)
                                remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, base_leave_type_name)
                                if not error and remaining:
                                    formatted = self.leave_balance_service.format_remaining_leave_message(remaining)
                                    if formatted:
                                        # Format as "**Available Annual Leave:** X days" or "**Available Sick Leave:** X days"
                                        import re
                                        formatted_bold_label = re.sub(r'(Available [^:]+):', r'**\1:**', formatted)
                                        # Use 💊 emoji for sick leave custom hours flow, ⏰ for others
                                        emoji = "💊" if base_leave_type_name == 'Sick Leave' else "⏰"
                                        remaining_leave_text = f"\n{emoji} {formatted_bold_label}"
                        except Exception as e:
                            debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                            # Continue without remaining leave info if there's an error
                    
                    # Check if this is a sick leave custom hours flow
                    is_sick_custom_hours = (
                        isinstance(selected_type, dict) and 
                        selected_type.get('name') == 'Sick Leave' and
                        self._is_custom_hours_mode(updated_session, 'Sick Leave')
                    )
                    
                    # Check if this is unpaid leave (unlimited, no balance to show)
                    is_unpaid_leave = (
                        isinstance(selected_type, dict) and 
                        selected_type.get('name') == 'Unpaid Leave'
                    )
                    
                    # Collect supporting documents if any
                    docs = self._collect_supporting_documents(updated_session)
                    doc_name = ''
                    if docs:
                        first_doc = docs[0]
                        doc_name = first_doc.get('filename') or first_doc.get('name') or 'Supporting document'
                    
                    if is_sick_custom_hours:
                        # Custom format for sick leave custom hours
                        response_text = f"Great, noted your hours. Here's your time-off request summary:\n\n"
                        response_text += f"📋 **Leave Type:** {leave_type_name}\n"
                        response_text += f"📅 **Date:** {dd_mm_yyyy(start_date)}\n"
                        response_text += f"⏰ **Hours:** {self._format_hour_label(raw_from)} to {self._format_hour_label(raw_to)}\n"
                        response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n"
                        if remaining_leave_text:
                            # Remove leading newline if present to avoid extra space
                            cleaned_remaining = remaining_leave_text.lstrip('\n')
                            # Ensure it ends with newline for proper formatting
                            if not cleaned_remaining.endswith('\n'):
                                cleaned_remaining += '\n'
                            response_text += cleaned_remaining
                        if doc_name:
                            response_text += f"📎 **Supporting Document:** {doc_name}\n"
                        response_text += "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    elif is_unpaid_leave:
                        # Format for unpaid leave (no leave balance - unlimited)
                        response_text = f"Great, noted your hours. Here's your time-off request summary:\n\n"
                        response_text += f"📋 **Leave Type:** {leave_type_name}\n"
                        response_text += f"📅 **Date:** {dd_mm_yyyy(start_date)}\n"
                        response_text += f"⏰ **Hours:** from {self._format_hour_label(raw_from)} to {self._format_hour_label(raw_to)}\n"
                        response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n"
                        response_text += "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    else:
                        # Original format for other leave types (with leave balance)
                        response_text = f"Great, noted your hours. Here's your time-off request summary:\n\n"
                        response_text += f"📋 **Leave Type:** {leave_type_name}\n"
                        response_text += f"📅 **Date:** {dd_mm_yyyy(start_date)}\n"
                        response_text += f"⏰ **Hours:** from {self._format_hour_label(raw_from)} to {self._format_hour_label(raw_to)}\n"
                        response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}{remaining_leave_text}\n\n"
                        response_text += "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    buttons = [
                        {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                        {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                    ]
                    return self._create_response_with_choice_buttons(response_text, thread_id, buttons)
            except Exception:
                pass

            # Single-date flow UX: if invalid hour range (e.g., same start/end) or unparsed text, re-open the hour picker
            try:
                if self._is_halfday_flow(session) or self._is_sick_custom_hours_flow(session):
                    hour_text = (
                        "Please choose a valid hours range (end must be after start)."
                    )
                    return self._create_response_with_hour_picker(hour_text, thread_id)
            except Exception:
                pass

            debug_log(f"Invalid confirmation response: '{message_lower}'", "bot_logic")
            # Instead of a bare prompt, re-show a full confirmation summary when possible
            try:
                session_data = session.get('data', {})
                context_data = {}
                if isinstance(session_data, dict):
                    context_data = session_data.get('timeoff_context', {}) or {}

                selected_type = (
                    session_data.get('selected_leave_type')
                    or session.get('selected_leave_type', {})
                    or context_data.get('selected_leave_type')
                    or {}
                )
                start_date = (
                    session_data.get('start_date')
                    or session.get('start_date')
                    or context_data.get('start_date')
                )
                end_date = (
                    session_data.get('end_date')
                    or session.get('end_date')
                    or context_data.get('end_date')
                )
                resolved_employee = (
                    employee_data
                    or session_data.get('employee_data')
                    or context_data.get('employee_data')
                    or session.get('employee_data')
                    or {}
                )

                if selected_type and start_date and end_date:
                    def dd_mm_yyyy(d: str) -> str:
                        try:
                            return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                        except Exception:
                            return d
                    
                    # Get remaining leave time (skip for unpaid leave - unlimited)
                    remaining_leave_text = ""
                    # Check if this is unpaid leave (unlimited, no balance to show)
                    is_unpaid_leave_check = (
                        isinstance(selected_type, dict) and 
                        selected_type.get('name') == 'Unpaid Leave'
                    )
                    
                    if not is_unpaid_leave_check:
                        try:
                            if self.leave_balance_service and resolved_employee and resolved_employee.get('id'):
                                employee_id = resolved_employee.get('id')
                                selected_leave_type_name = selected_type.get('name', 'Unknown')
                                remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, selected_leave_type_name)
                                if not error and remaining:
                                    formatted = self.leave_balance_service.format_remaining_leave_message(remaining)
                                    if formatted:
                                        # Format as "**Available Annual Leave:** X days" (only label bold, not days)
                                        import re
                                        formatted_bold_label = re.sub(r'(Available [^:]+):', r'**\1:**', formatted)
                                        remaining_leave_text = f"\n⏰ {formatted_bold_label}"
                        except Exception as e:
                            debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                            # Continue without remaining leave info if there's an error
                    
                    response_text = f"Great, noted your dates. Here's your time-off request summary:\n\n"
                    response_text += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown')}\n"
                    response_text += f"📅 **Start Date:** {dd_mm_yyyy(start_date)}\n"
                    response_text += f"📅 **End Date:** {dd_mm_yyyy(end_date)}\n"
                    response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}{remaining_leave_text}\n"
                    response_text += "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    buttons = [
                        {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                        {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                    ]
                    return self._create_response_with_choice_buttons(response_text, thread_id, buttons)

                # If we don't have enough context, direct the user to the next required step
                if not selected_type:
                    leave_types = session.get('main_leave_types') or session.get('leave_types', [])
                    if leave_types:
                        prompt = "I lost track of which leave type you picked. Please choose it again:"
                        return self._create_response_with_buttons(prompt, thread_id, leave_types)
                    return self._create_response("I lost track of which leave type you picked. Please tell me the leave type.", thread_id)

                reprompt = (
                    "I still need both the start and end date before I can submit the request. "
                    "Please send them in one message, for example '15/10/2025 to 16/10/2025'."
                )
                return self._create_response_with_datepicker(reprompt, thread_id)
            except Exception:
                # Fallback minimal prompt only if summary rendering fails for some reason
                response_text = "Please reply with 'yes' to submit the request or 'no' to cancel."
                return self._create_response(response_text, thread_id)
    
    def _submit_timeoff_request(self, thread_id: str, session: dict, employee_data: dict, odoo_session_data: dict = None) -> dict:
        """Submit the time-off request to Odoo (with optional session data for stateless mode)"""
        try:
            debug_log(f"Starting time-off submission for thread: {thread_id}", "bot_logic")
            debug_log(f"Full session structure: {session}", "bot_logic")

            session_data = session.get('data', {})
            context_data = {}
            if isinstance(session_data, dict):
                context_data = session_data.get('timeoff_context', {}) or {}

            # Check both session root level and nested data level
            selected_type = (
                session_data.get('selected_leave_type')
                or session.get('selected_leave_type', {})
                or context_data.get('selected_leave_type')
                or {}
            )
            start_date = (
                session_data.get('start_date')
                or session.get('start_date')
                or context_data.get('start_date')
            )
            end_date = (
                session_data.get('end_date')
                or session.get('end_date')
                or context_data.get('end_date')
            )
            # Resolve employee data from argument or session to ensure correctness
            resolved_employee = (
                employee_data
                or session_data.get('employee_data')
                or context_data.get('employee_data')
                or session.get('employee_data')
                or {}
            )

            debug_log(f"Session data - selected_type: {selected_type}, start_date: {start_date}, end_date: {end_date}", "bot_logic")

            employee_id = resolved_employee.get('id') if resolved_employee else None
            leave_type_id = selected_type.get('id') if selected_type else None

            debug_log(f"Extracted data - employee_id: {employee_id}, leave_type_id: {leave_type_id}", "bot_logic")

            if not all([employee_id, leave_type_id, start_date, end_date]):
                debug_log(f"Missing required data - employee_id: {employee_id}, leave_type_id: {leave_type_id}, start_date: {start_date}, end_date: {end_date}", "general")
                raise ValueError("Missing required data for submission")
            
            # Support Half Day custom hours and Sick Leave Custom Hours via modular service and hour range fields
            extra_fields = {}
            supporting_documents = self._collect_supporting_documents(session)

            if self._requires_supporting_document(session) and not supporting_documents:
                debug_log("Supporting document required but not available; prompting upload.", "bot_logic")
                reminder_response = self._prompt_supporting_document_upload(thread_id)
                reminder_response['message'] = (
                    "I still need your supporting document before I can submit this Sick Leave request. "
                    "Please upload it using the button above."
                )
                return reminder_response
            
            # Check if this is a Custom Hours mode for the current leave type
            selected_name_sub = ''
            is_custom_hours_mode = False
            try:
                selected_name_sub = (selected_type.get('name') or '').strip()
                is_custom_hours_mode = self._is_custom_hours_mode(session, selected_name_sub)
            except Exception:
                is_custom_hours_mode = False
            
            # Handle halfday service or custom-hours flows
            if (self.halfday_service and isinstance(selected_type, dict)) or is_custom_hours_mode:
                # Only map leave type ID for halfday service (not for Sick/Unpaid Custom Hours)
                if self.halfday_service and isinstance(selected_type, dict) and not is_custom_hours_mode:
                    try:
                        mapped_leave_type_id, hd_extra = self.halfday_service.build_submission(selected_type)
                        if mapped_leave_type_id:
                            leave_type_id = mapped_leave_type_id
                        if hd_extra:
                            extra_fields.update(hd_extra)
                    except Exception as hd_map_e:
                        debug_log(f"HalfDay build_submission error: {hd_map_e}", "general")
                
                # For Sick/Unpaid Custom Hours, enable the Custom Hours option in Odoo
                if is_custom_hours_mode:
                    extra_fields['request_unit_hours'] = True
                    debug_log(f"{selected_name_sub or 'Custom'} Custom Hours: setting request_unit_hours=True", "bot_logic")
                
                # Add hour range if present in session (for both halfday and Sick Leave Custom Hours)
                sd = session.get('data', {})
                hour_from = sd.get('hour_from') or session.get('hour_from') or context_data.get('hour_from')
                hour_to = sd.get('hour_to') or session.get('hour_to') or context_data.get('hour_to')

                # Enforce 4-hour maximum
                try:
                    hf = float(hour_from) if hour_from is not None else float('nan')
                    ht = float(hour_to) if hour_to is not None else float('nan')
                    if hf == hf and ht == ht and (ht - hf) > 4.0 - 1e-9:
                        return self._create_response_with_hour_picker(
                            "Per Prezlab policy, the maximum Half Day leave duration is 4 hours. Please choose up to 4 hours.",
                            thread_id
                        )
                except Exception:
                    pass

                def _to_selection_key(val):
                    # Convert various numeric representations to the canonical Odoo selection format (e.g. '9', '9.5')
                    try:
                        f = float(val) if isinstance(val, (str, int, float)) else None
                        if f is None:
                            return None
                        # Round to nearest half-hour just in case (align with picker granularity)
                        f = round(f * 2) / 2.0
                        if abs(f - round(f)) < 1e-9:
                            return str(int(round(f)))
                        return f"{f:.1f}".rstrip('0').rstrip('.')
                    except Exception:
                        return None

                hf = _to_selection_key(hour_from)
                ht = _to_selection_key(hour_to)
                if hf is not None:
                    extra_fields['request_hour_from'] = hf
                if ht is not None:
                    extra_fields['request_hour_to'] = ht
                    try:
                        fhf = float(hf) if hf is not None else float('nan')
                        fht = float(ht)
                        if fhf == fhf and fht == fht and fht > fhf:
                            duration_hours = round(fht - fhf, 2)
                            extra_fields.setdefault('number_of_hours_display', duration_hours)
                            extra_fields.setdefault('number_of_hours', duration_hours)
                    except Exception:
                        pass

            attachment_payloads: List[Dict[str, Any]] = []
            for doc in supporting_documents:
                if not isinstance(doc, dict):
                    continue
                datas = doc.get('data')
                if not datas:
                    continue
                attachment_payloads.append({
                    'filename': doc.get('filename') or doc.get('name'),
                    'mimetype': doc.get('mimetype') or doc.get('content_type'),
                    'data': datas
                })

            # Use stateless version if session data provided, otherwise fall back to old method
            if odoo_session_data and odoo_session_data.get('session_id') and odoo_session_data.get('user_id'):
                print(f"DEBUG: Using stateless submit with session_id: {odoo_session_data.get('session_id')[:20]}...")
                success, result, renewed_session = self.timeoff_service.submit_leave_request_stateless(
                    employee_id=employee_id,
                    leave_type_id=leave_type_id,
                    start_date=start_date,
                    end_date=end_date,
                    description=f"Time off request submitted via Nasma chatbot",
                    extra_fields=extra_fields or None,
                    supporting_attachments=attachment_payloads or None,
                    session_id=odoo_session_data.get('session_id'),
                    user_id=odoo_session_data.get('user_id'),
                    username=odoo_session_data.get('username'),
                    password=odoo_session_data.get('password')
                )

                # CRITICAL: If session was renewed, update Flask session
                if renewed_session:
                    print(f"DEBUG: Updating Flask session with renewed Odoo session_id")
                    try:
                        from flask import session as flask_session
                        flask_session['odoo_session_id'] = renewed_session['session_id']
                        flask_session['user_id'] = renewed_session['user_id']
                        flask_session.modified = True
                        print(f"DEBUG: Flask session updated successfully with new session_id: {renewed_session['session_id'][:20]}...")
                    except Exception as session_update_error:
                        print(f"DEBUG: Failed to update Flask session: {session_update_error}")
            else:
                print("DEBUG: Using legacy submit (no session data provided)")
                success, result = self.timeoff_service.submit_leave_request(
                    employee_id=employee_id,
                    leave_type_id=leave_type_id,
                    start_date=start_date,
                    end_date=end_date,
                    description=f"Time off request submitted via Nasma chatbot",
                    extra_fields=extra_fields or None,
                    supporting_attachments=attachment_payloads or None
                )
            
            if success:
                extras = extra_fields if isinstance(extra_fields, dict) else {}
                metric_payload = {
                    'leave_id': result.get('leave_id') if isinstance(result, dict) else None,
                    'employee_id': employee_id,
                    'leave_type_id': leave_type_id,
                    'leave_type_name': (
                        selected_type.get('name') if isinstance(selected_type, dict) else None
                    ),
                    'start_date': start_date,
                    'end_date': end_date,
                    'hours': {
                        'request_unit_hours': bool(extras.get('request_unit_hours')) if 'request_unit_hours' in extras else None,
                        'number_of_hours': extras.get('number_of_hours'),
                        'number_of_hours_display': extras.get('number_of_hours_display'),
                    },
                    'attachments_submitted': len(attachment_payloads),
                }
                if not self._record_metric('timeoff', thread_id, metric_payload, resolved_employee or employee_data):
                    msg = getattr(self.metrics_service, "last_error", None) if self.metrics_service else None
                    debug_log(f"Supabase metric logging failed for timeoff: {msg}", "general")
                self.session_manager.complete_session(thread_id, {'submitted': True, 'result': result})
                try:
                    self._reset_timeoff_flow_state(thread_id)
                except Exception:
                    pass
                response_text = f"✅ **Success!** {result.get('message', 'Your time-off request has been submitted.')}\n\n"
                response_text += "Your request is now pending approval from your manager. You should receive a notification once it's reviewed."
            else:
                self.session_manager.complete_session(thread_id, {'submitted': False, 'error': result})
                try:
                    self._reset_timeoff_flow_state(thread_id)
                except Exception:
                    pass
                response_text = f"❌ **Submission Failed:** {result}\n\n"
                response_text += "Please try again later or contact your HR department for assistance."
            
            resp = self._create_response(response_text, thread_id)
            try:
                # Mark as handled to avoid any lingering prompts
                resp['session_handled'] = True
            except Exception:
                pass
            return resp
            
        except Exception as e:
            debug_log(f"Error submitting time-off request: {e}", "general")
            self.session_manager.complete_session(thread_id, {'submitted': False, 'error': str(e)})
            response_text = "❌ Sorry, there was an error submitting your request. Please try again later or contact HR for assistance."
            return self._create_response(response_text, thread_id)
    
    def _create_response(self, message_text: str, thread_id: str) -> dict:
        """Create a standard response object"""
        debug_log(f"Creating response with message length: {len(message_text) if message_text else 0}", "bot_logic")
        debug_log(f"Response thread_id: {thread_id}", "bot_logic")

        # Ensure we always have a thread_id for session continuity
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response: {thread_id}", "bot_logic")

        result = {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True
        }

        debug_log(f"Response object created successfully: {result is not None}", "bot_logic")
        return result

    def _create_response_with_buttons(self, message_text: str, thread_id: str, leave_types: list) -> dict:
        """Create a response object with interactive buttons"""
        debug_log(f"Creating response with buttons - message length: {len(message_text) if message_text else 0}", "bot_logic")
        debug_log(f"Response thread_id: {thread_id}", "bot_logic")

        # Ensure we always have a thread_id for session continuity
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with buttons: {thread_id}", "bot_logic")

        # Create button data
        buttons = []
        for leave_type in leave_types:
            buttons.append({
                'text': leave_type.get('name', 'Unknown'),
                'value': leave_type.get('name', 'Unknown'),
                'type': 'leave_type_selection'
            })

        result = {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'buttons': buttons  # Add buttons to response
        }

        debug_log(f"Response object with {len(buttons)} buttons created successfully", "bot_logic")
        return result

    def _create_response_with_choice_buttons(self, message_text: str, thread_id: str,
                                             buttons: list) -> dict:
        """Create a response with explicit buttons, each a dict with text/value/type."""
        debug_log(f"Creating response with choice buttons - count: {len(buttons) if buttons else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'buttons': buttons
        }

    def _prompt_sick_leave_mode(self, thread_id: str, prefix_note: str = "", leave_name: str = "Sick Leave") -> dict:
        """Backward-compatible alias that delegates to the generic leave-mode prompt."""
        return self._prompt_leave_mode(thread_id, leave_name, prefix_note)

    def _requires_supporting_document(self, session: dict) -> bool:
        """Return True when current selection requires a supporting document upload."""
        try:
            session_data = session.get('data', {}) if isinstance(session, dict) else {}
            ctx = {}
            if isinstance(session_data, dict):
                ctx = session_data.get('timeoff_context', {}) or {}
            selected = (
                ctx.get('selected_leave_type')
                or session_data.get('selected_leave_type')
                or session.get('selected_leave_type')
                or {}
            )
            name = (selected.get('name') or '').strip().lower()
            if name != 'sick leave':
                return False
            mode = self._get_leave_mode(session, 'Sick Leave')
            if isinstance(mode, str):
                mode = mode.strip().lower()
            return mode in {
                'full_days', 'full day', 'fulldays', 'full-days',
                'custom_hours', 'custom hour', 'custom-hours'
            }
        except Exception:
            return False

    def _collect_supporting_documents(self, session: dict) -> List[dict]:
        """Return supporting documents (if any) stored in the session."""
        try:
            session_data = session.get('data', {}) if isinstance(session, dict) else {}
            docs = session_data.get('supporting_documents')
            if docs:
                return docs if isinstance(docs, list) else []
            docs = session.get('supporting_documents')
            if docs and isinstance(docs, list):
                return docs
            ctx = session_data.get('timeoff_context', {}) if isinstance(session_data, dict) else {}
            docs = ctx.get('supporting_documents')
            return docs if isinstance(docs, list) else []
        except Exception:
            return []

    def _has_supporting_document(self, session: dict) -> bool:
        """Check whether the session already has an uploaded supporting document."""
        docs = self._collect_supporting_documents(session)
        return any(isinstance(doc, dict) and doc.get('data') for doc in docs)

    def _is_supporting_doc_stage(self, session: dict) -> bool:
        """Determine if we are waiting for a supporting document upload."""
        try:
            session_data = session.get('data', {}) if isinstance(session, dict) else {}
            required = (
                session_data.get('supporting_doc_required')
                if isinstance(session_data, dict) else False
            )
            uploaded = (
                session_data.get('supporting_doc_uploaded')
                if isinstance(session_data, dict) else False
            )
            if not required:
                required = session.get('supporting_doc_required', False)
            if not uploaded:
                uploaded = session.get('supporting_doc_uploaded', False)
            return bool(required) and not bool(uploaded)
        except Exception:
            return False

    def _prompt_supporting_document_upload(self, thread_id: str) -> dict:
        """Create a response requesting the supporting document upload widget."""
        message_text = (
            "Sick Leave full days require a medical certificate or supporting document. "
            "Please upload the document now so I can submit your request."
        )
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'widgets': {
                'supporting_document_upload': {
                    'context_key': 'timeoff_support_document',
                    'accept': '.pdf,.jpg,.jpeg,.png,.heic,.doc,.docx',
                    'max_files': 1
                }
            }
        }

    def _handle_supporting_document_step(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle the supporting document upload step before confirmation."""
        message_lower = (message or '').strip().lower()
        import re as _re_support
        message_clean = _re_support.sub(r'[^\w\s]', '', message_lower)
        cancel_tokens = {'cancel', 'abort', 'stop', 'exit', 'quit', 'no', 'n', 'nevermind', 'undo', 'end'}
        if message_lower in cancel_tokens or message_clean in cancel_tokens:
            try:
                self.session_manager.cancel_session(thread_id, 'User cancelled during supporting document step')
            finally:
                self._reset_timeoff_flow_state(thread_id)
            return self._create_response('request cancelled, can i help you with anything else', thread_id)

        # Allow restart commands mid-step
        if self._is_timeoff_start_message(message):
            return self._restart_timeoff_flow(message, thread_id, employee_data, 'User restarted during supporting document step')

        # Accept explicit acknowledgement tokens after upload
        acknowledgement_tokens = {
            'supporting_document_uploaded',
            'supporting document uploaded',
            'document uploaded',
            'uploaded',
            'done',
            'finished'
        }
        session_refreshed = session
        try:
            refreshed = self.session_manager.get_session(thread_id)
            if refreshed:
                session_refreshed = refreshed
        except Exception:
            pass

        has_doc = self._has_supporting_document(session_refreshed)
        if message_lower in acknowledgement_tokens or message_clean in acknowledgement_tokens:
            if not has_doc:
                reminder = (
                    "I still don't see a document attached. Please use the upload button above to share your medical certificate."
                )
                return self._prompt_supporting_document_upload(thread_id) | {'message': reminder}

            # Mark the supporting document as uploaded and proceed to confirmation
            try:
                current_step = session_refreshed.get('step', 2)
                updates = {
                    'supporting_doc_uploaded': True,
                }
                self.session_manager.update_session(thread_id, updates)
                # Store in data for convenience
                data_payload = session_refreshed.get('data', {}) or {}
                ctx_payload = dict((data_payload.get('timeoff_context') or {}))
                ctx_payload.update({
                    'supporting_doc_uploaded': True,
                    'supporting_doc_required': True,
                })
                data_payload.update({
                    'supporting_doc_uploaded': True,
                    'supporting_doc_required': True,
                    'timeoff_context': ctx_payload
                })
                self.session_manager.update_session(thread_id, {'data': data_payload})
                # Advance to the confirmation step now that docs are ready
                if current_step <= 2:
                    self.session_manager.advance_session_step(thread_id)
            except Exception:
                pass

            # Build the confirmation summary (same as after dates)
            ctx = self._resolve_timeoff_context(session_refreshed)
            selected_type = ctx.get('selected_leave_type') or session_refreshed.get('selected_leave_type') or {}
            start_date = ctx.get('start_date') or session_refreshed.get('start_date')
            end_date = ctx.get('end_date') or session_refreshed.get('end_date')
            resolved_employee = (
                employee_data
                or ctx.get('employee_data')
                or session_refreshed.get('employee_data')
                or {}
            )

            def dd_mm_yyyy(d: str) -> str:
                try:
                    return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                except Exception:
                    return d

            docs = self._collect_supporting_documents(session_refreshed)
            doc_name = ''
            if docs:
                first_doc = docs[0]
                doc_name = first_doc.get('filename') or first_doc.get('name') or 'Supporting document'

            # Check if this is a sick leave custom hours flow
            is_sick_custom_hours = (
                isinstance(selected_type, dict) and 
                selected_type.get('name') == 'Sick Leave' and
                self._is_custom_hours_mode(session_refreshed, 'Sick Leave')
            )
            
            # Check if this is unpaid leave (unlimited, no balance to show)
            is_unpaid_leave = (
                isinstance(selected_type, dict) and 
                selected_type.get('name') == 'Unpaid Leave'
            )
            
            # Get hours for custom hours flow
            raw_from = ctx.get('hour_from') or session_refreshed.get('hour_from')
            raw_to = ctx.get('hour_to') or session_refreshed.get('hour_to')
            
            # Get remaining leave time (skip for unpaid leave - unlimited)
            remaining_leave_text = ""
            if not is_unpaid_leave:
                try:
                    if self.leave_balance_service and resolved_employee and resolved_employee.get('id'):
                        employee_id = resolved_employee.get('id')
                        # Determine base leave type name for balance calculation
                        base_leave_type_name = selected_type.get('name', 'Unknown') if isinstance(selected_type, dict) else 'Unknown'
                        
                        # Check if this is Half Days (Custom Hours) flow - should use Annual Leave for balance
                        try:
                            if self.halfday_service and isinstance(selected_type, dict):
                                if self.halfday_service.is_halfday(selected_type):
                                    # For Half Days, use Annual Leave for balance calculation
                                    base_leave_type_name = 'Annual Leave'
                        except Exception:
                            pass
                        
                        remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, base_leave_type_name)
                        if not error and remaining:
                            formatted = self.leave_balance_service.format_remaining_leave_message(remaining)
                            if formatted:
                                # Format as "**Available Annual Leave:** X days" or "**Available Sick Leave:** X days"
                                import re
                                formatted_bold_label = re.sub(r'(Available [^:]+):', r'**\1:**', formatted)
                                # Use 💊 emoji for sick leave custom hours flow, ⏰ for others
                                emoji = "💊" if base_leave_type_name == 'Sick Leave' else "⏰"
                                remaining_leave_text = f"\n{emoji} {formatted_bold_label}"
                except Exception as e:
                    debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                    # Continue without remaining leave info if there's an error

            if is_sick_custom_hours and raw_from and raw_to:
                # Custom format for sick leave custom hours
                summary = "Great, I have everything I need. Here's your time-off request summary:\n\n"
                summary += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if isinstance(selected_type, dict) else 'Unknown'}\n"
                summary += f"📅 **Date:** {dd_mm_yyyy(start_date) if start_date else 'Unknown'}\n"
                summary += f"⏰ **Hours:** {self._format_hour_label(raw_from)} to {self._format_hour_label(raw_to)}\n"
                summary += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n"
                if remaining_leave_text:
                    # Remove leading newline if present to avoid extra space
                    cleaned_remaining = remaining_leave_text.lstrip('\n')
                    # Ensure it ends with newline for proper formatting
                    if not cleaned_remaining.endswith('\n'):
                        cleaned_remaining += '\n'
                    summary += cleaned_remaining
                if doc_name:
                    summary += f"📎 **Supporting Document:** {doc_name}\n"
                summary += "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
            elif is_unpaid_leave and raw_from and raw_to:
                # Format for unpaid leave custom hours (no leave balance - unlimited)
                summary = "Great, I have everything I need. Here's your time-off request summary:\n\n"
                summary += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if isinstance(selected_type, dict) else 'Unknown'}\n"
                summary += f"📅 **Date:** {dd_mm_yyyy(start_date) if start_date else 'Unknown'}\n"
                summary += f"⏰ **Hours:** {self._format_hour_label(raw_from)} to {self._format_hour_label(raw_to)}\n"
                summary += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n"
                if doc_name:
                    summary += f"📎 **Supporting Document:** {doc_name}\n"
                summary += "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
            else:
                # Original format for other leave types
                summary = "Great, I have everything I need. Here's your time-off request summary:\n\n"
                summary += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if isinstance(selected_type, dict) else 'Unknown'}\n"
                summary += f"📅 **Start Date:** {dd_mm_yyyy(start_date) if start_date else 'Unknown'}\n"
                summary += f"📅 **End Date:** {dd_mm_yyyy(end_date) if end_date else 'Unknown'}\n"
                summary += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}{remaining_leave_text}\n"
                if doc_name:
                    summary += f"📎 **Supporting Document:** {doc_name}\n"
                summary += "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
            buttons = [
                {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
            ]
            return self._create_response_with_choice_buttons(summary, thread_id, buttons)

        # Any other message should re-prompt the upload instructions
        reminder = (
            "Please upload your supporting document using the button above. "
            "Once it's uploaded, click the Upload button and I'll continue."
        )
        return self._prompt_supporting_document_upload(thread_id) | {'message': reminder}

    def _process_captured_dates(self, thread_id: str, session: dict, employee_data: dict,
                                start_date: str, end_date: str, is_single_date_flow: bool) -> dict:
        """Store captured dates then branch to supporting document prompt or confirmation summary."""
        self.session_manager.update_session(thread_id, {'start_date': start_date, 'end_date': end_date})
        self._persist_timeoff_context(thread_id, session, start_date=start_date, end_date=end_date)

        refreshed_session = session
        try:
            latest = self.session_manager.get_session(thread_id)
            if latest:
                refreshed_session = latest
        except Exception:
            pass

        if self._requires_supporting_document(refreshed_session) and not is_single_date_flow:
            try:
                data_payload = dict(refreshed_session.get('data', {}) or {})
                ctx = dict(data_payload.get('timeoff_context', {}) or {})
                ctx.update({
                    'start_date': start_date,
                    'end_date': end_date,
                    'supporting_doc_required': True,
                })
                data_payload['timeoff_context'] = ctx
                data_payload['supporting_doc_required'] = True
                data_payload['supporting_doc_uploaded'] = False
                updates = {
                    'supporting_doc_required': True,
                    'supporting_doc_uploaded': False,
                    'data': data_payload
                }
                self.session_manager.update_session(thread_id, updates)
            except Exception:
                pass
            return self._prompt_supporting_document_upload(thread_id)

        # No supporting document requirement; proceed normally
        self.session_manager.advance_session_step(thread_id)
        try:
            self.session_manager.update_session(thread_id, {'step': 3})
        except Exception:
            pass

        ctx = self._resolve_timeoff_context(refreshed_session)
        selected_type = ctx.get('selected_leave_type') or refreshed_session.get('selected_leave_type') or {}
        resolved_employee = (
            employee_data
            or ctx.get('employee_data')
            or refreshed_session.get('employee_data')
            or {}
        )

        if is_single_date_flow:
            return self._create_response_with_hour_picker(
                "Great, got your date. Please choose your hours (from/to).",
                thread_id
            )

        def dd_mm_yyyy(d: str) -> str:
            try:
                return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
            except Exception:
                return d

        # Get remaining leave time (skip for unpaid leave - unlimited)
        remaining_leave_text = ""
        # Check if this is unpaid leave (unlimited, no balance to show)
        is_unpaid_leave_check = (
            isinstance(selected_type, dict) and 
            selected_type.get('name') == 'Unpaid Leave'
        )
        
        if not is_unpaid_leave_check:
            try:
                if self.leave_balance_service and resolved_employee and resolved_employee.get('id'):
                    employee_id = resolved_employee.get('id')
                    leave_type_name = selected_type.get('name', 'Unknown') if isinstance(selected_type, dict) else 'Unknown'
                    remaining, error = self.leave_balance_service.calculate_remaining_leave(employee_id, leave_type_name)
                    if not error and remaining:
                        formatted = self.leave_balance_service.format_remaining_leave_message(remaining)
                        if formatted:
                            # Format as "**Available Annual Leave:** X days" (only label bold, not days)
                            import re
                            formatted_bold_label = re.sub(r'(Available [^:]+):', r'**\1:**', formatted)
                            remaining_leave_text = f"\n⏰ {formatted_bold_label}"
            except Exception as e:
                debug_log(f"Error getting remaining leave: {str(e)}", "bot_logic")
                # Continue without remaining leave info if there's an error

        summary = "Great, noted your dates. Here's your time-off request summary:\n\n"
        summary += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if isinstance(selected_type, dict) else 'Unknown'}\n"
        summary += f"📅 **Start Date:** {dd_mm_yyyy(start_date)}\n"
        summary += f"📅 **End Date:** {dd_mm_yyyy(end_date)}\n"
        summary += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}{remaining_leave_text}\n"
        summary += "\nDo you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
        buttons = [
            {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
            {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
        ]
        return self._create_response_with_choice_buttons(summary, thread_id, buttons)

    def _create_response_with_datepicker(self, message_text: str, thread_id: str) -> dict:
        """Create a response object that instructs the UI to show a date range picker widget"""
        debug_log(f"Creating response with datepicker - message length: {len(message_text) if message_text else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with datepicker: {thread_id}", "bot_logic")

        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'widgets': {
                'date_range_picker': True
            }
        }

    def _create_response_with_datepicker_single(self, message_text: str, thread_id: str) -> dict:
        """Create a response object to show a single-date picker widget (for Half Days)."""
        debug_log(f"Creating response with single-date picker - message length: {len(message_text) if message_text else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with single datepicker: {thread_id}", "bot_logic")

        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'widgets': {
                'single_date_picker': True
            }
        }

    def _create_response_with_hour_picker(self, message_text: str, thread_id: str) -> dict:
        """Create a response object that instructs the UI to show an hour-from/to picker."""
        debug_log(f"Creating response with hour picker - message length: {len(message_text) if message_text else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with hour picker: {thread_id}", "bot_logic")

        # Build hour options (0..23.5 step 0.5)
        def _fmt_label(val: float) -> str:
            h = int(val)
            m = 30 if abs(val - h - 0.5) < 1e-6 else 0
            ampm_h = h % 12
            if ampm_h == 0:
                ampm_h = 12
            ampm = 'AM' if h < 12 else 'PM'
            return f"{ampm_h}:{m:02d} {ampm}" if m else f"{ampm_h}:00 {ampm}"

        options = []
        def _push(val: float):
            canonical = round(val * 2) / 2.0
            key = f"{canonical:.1f}"
            options.append({'value': key, 'label': _fmt_label(val)})

        # From 9:00 to 23:30 (same day)
        v = 9.0
        while v <= 23.5 + 1e-9:
            _push(v)
            v += 0.5
        # Then wrap 0:00 to 1:00 (next day)
        v = 0.0
        while v <= 1.0 + 1e-9:
            _push(v)
            v += 0.5

        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'widgets': {
                'hour_range_picker': True,
                'hour_options': options
            }
        }

    def _is_halfday_flow(self, session: dict) -> bool:
        session_data = session.get('data', {}) if isinstance(session, dict) else {}
        selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
        try:
            if self.halfday_service and isinstance(selected_type, dict):
                return self.halfday_service.is_halfday(selected_type) or (
                    selected_type.get('name') == getattr(self.halfday_service, 'HALF_DAY_NAME', 'Custom Hours')
                )
        except Exception:
            return False
        return False
    
    def _is_sick_custom_hours_flow(self, session: dict) -> bool:
        """Check if current flow is a custom-hours mode (Sick/Unpaid)."""
        session_data = session.get('data', {}) if isinstance(session, dict) else {}
        selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
        try:
            selected_name = (selected_type.get('name') or '').strip()
            return self._is_custom_hours_mode(session, selected_name)
        except Exception:
            return False

    def _format_hour_label(self, value: str) -> str:
        """Format a float-like string hour value to 12-hour label (e.g., '8.5' -> '8:30 AM')."""
        try:
            f = float(value)
            h = int(f)
            minutes = 30 if abs(f - h - 0.5) < 1e-6 else 0
            ampm_h = h % 12
            if ampm_h == 0:
                ampm_h = 12
            ampm = 'AM' if h < 12 else 'PM'
            return f"{ampm_h}:{minutes:02d} {ampm}" if minutes else f"{ampm_h}:00 {ampm}"
        except Exception:
            return value

    def _create_response_with_choice_buttons(self, message_text: str, thread_id: str,
                                             buttons: list) -> dict:
        """Create a response with explicit buttons, each a dict with text/value/type."""
        debug_log(f"Creating response with choice buttons - count: {len(buttons) if buttons else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'buttons': buttons
        }

    def _parse_hour_value(self, token: str) -> float:
        """Parse a time token like '4pm', '16', '4:30 pm', '16:30' to 24h float (e.g., 16.5)."""
        try:
            import re
            s = token.strip().lower()
            s = s.replace('.', ':')
            m = re.match(r"^(\d{1,2})(?::(\d{2}))?\s*(am|pm)?$", s)
            if not m:
                return float('nan')
            h = int(m.group(1))
            mins = int(m.group(2)) if m.group(2) else 0
            ap = m.group(3)
            if ap == 'am':
                if h == 12:
                    h = 0
            elif ap == 'pm':
                if h != 12:
                    h += 12
            # clamp to 0..23:59
            h = max(0, min(23, h))
            mins = 0 if mins < 15 else (30 if mins < 45 else 0)  # snap to nearest half hour (0 or 30)
            return h + (0.5 if mins == 30 else 0.0)
        except Exception:
            return float('nan')

    def _parse_hour_range_text(self, text: str):
        """Parse expressions like '4pm to 5pm', '16:00 - 17:30'. Return (from_str, to_str) or (None, None)."""
        try:
            import re
            if not text or not isinstance(text, str):
                return None, None
            s = text.strip().lower()
            # Normalize connectors
            s = re.sub(r"\s*(?:–|—)\s*", '-', s)
            parts = re.split(r"\s*(?:to|till|until|-)\s*", s)
            if len(parts) != 2:
                return None, None
            v1 = self._parse_hour_value(parts[0])
            v2 = self._parse_hour_value(parts[1])
            if v1 != v1 or v2 != v2:  # NaN check
                return None, None
            # Ensure within 0..24 and logical order
            if v2 <= v1:
                return None, None
            return f"{v1:.1f}", f"{v2:.1f}"
        except Exception:
            return None, None

    def _hour_key_to_float(self, val: str) -> float:
        try:
            return float(val)
        except Exception:
            return float('nan')

    # -------------------- Robust context helpers (time-off) --------------------
    def _resolve_timeoff_context(self, session: dict) -> dict:
        """Resolve time-off context fields from both root and nested session data.
        Returns a dict with keys: selected_leave_type, start_date, end_date, employee_data.
        """
        try:
            sd = session.get('data', {}) if isinstance(session, dict) else {}
            ctx = sd.get('timeoff_context', {}) if isinstance(sd, dict) else {}
            resolved = {
                'selected_leave_type': (
                    session.get('selected_leave_type')
                    or sd.get('selected_leave_type')
                    or ctx.get('selected_leave_type')
                ),
                'start_date': (
                    session.get('start_date')
                    or sd.get('start_date')
                    or ctx.get('start_date')
                ),
                'end_date': (
                    session.get('end_date')
                    or sd.get('end_date')
                    or ctx.get('end_date')
                ),
                'employee_data': (
                    sd.get('employee_data')
                    or ctx.get('employee_data')
                    or session.get('employee_data')
                    or {}
                ),
                'sick_leave_mode': (
                    session.get('sick_leave_mode')
                    or sd.get('sick_leave_mode')
                    or ctx.get('sick_leave_mode')
                ),
                'unpaid_leave_mode': (
                    session.get('unpaid_leave_mode')
                    or sd.get('unpaid_leave_mode')
                    or ctx.get('unpaid_leave_mode')
                )
            }
            return resolved
        except Exception:
            return {
                'selected_leave_type': None,
                'start_date': None,
                'end_date': None,
                'employee_data': {}
            }

    # -------------------- Full reset helpers --------------------
    def _clear_summary(self, thread_id: str) -> None:
        try:
            path = self._get_summary_file(thread_id)
            if os.path.exists(path):
                os.remove(path)
        except Exception:
            pass

    def _reset_timeoff_flow_state(self, thread_id: str) -> None:
        """Comprehensively reset any cached state for this thread after finish/cancel."""
        try:
            if thread_id:
                try:
                    self.session_manager.clear_session(thread_id)
                except Exception:
                    pass
                try:
                    self.clear_conversation_history(thread_id)
                except Exception:
                    pass
                try:
                    self._clear_summary(thread_id)
                except Exception:
                    pass
        except Exception:
            pass
